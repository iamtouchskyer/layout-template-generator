#!/usr/bin/env python3
"""
PPTX Generator - Generate PowerPoint from Layout Generator JSON config
Uses python-pptx Slide Master API for proper inheritance.

Usage: python pptx_generator.py config.json output.pptx
"""

import sys
import json
from pathlib import Path

# Add python-pptx to path
sys.path.insert(0, '/Users/touichskyer/Code/python-pptx/src')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.smartart import SMARTART_TYPE, SMARTART_COLORS
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.smartart import SmartArtData


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


# Theme color mappings
THEME_COLORS = {
    'soft_peach_cream': {
        'primary': '#3D3D3D', 'accent': '#F5A89A', 'bg': '#FFFAF5',
        'text': '#3D3D3D', 'text_muted': '#888888', 'card_bg': '#FFF5F0', 'card_border': '#F5D5C8'
    },
    'executive': {
        'primary': '#FFFFFF', 'accent': '#64B5F6', 'bg': '#1A2942',
        'text': '#FFFFFF', 'text_muted': '#8899AA', 'card_bg': '#243B55', 'card_border': '#3D5A80'
    },
    'forest_green': {
        'primary': '#1A3D2B', 'accent': '#2D5A3D', 'bg': '#F5F7F5',
        'text': '#1A3D2B', 'text_muted': '#6B8068', 'card_bg': '#E8F0E8', 'card_border': '#C8D8C8'
    },
    'sunset_orange': {
        'primary': '#2D2D2D', 'accent': '#E85D3B', 'bg': '#FFF9F5',
        'text': '#2D2D2D', 'text_muted': '#888888', 'card_bg': '#FFF0E8', 'card_border': '#F5D0C0'
    },
    'cosmic': {
        'primary': '#F2F2F3', 'accent': '#9B59B6', 'bg': '#050505',
        'text': '#F2F2F3', 'text_muted': '#888888', 'card_bg': '#1A1A1A', 'card_border': '#333333'
    },
    'azure': {
        'primary': '#002FA7', 'accent': '#002FA7', 'bg': '#F8FAFC',
        'text': '#1A1A2E', 'text_muted': '#6B7280', 'card_bg': '#EEF2FF', 'card_border': '#C7D2FE'
    },
    'bright_red_blue': {
        'primary': '#101010', 'accent': '#E4380E', 'bg': '#F2F5F6',
        'text': '#101010', 'text_muted': '#666666', 'card_bg': '#FFFFFF', 'card_border': '#E0E0E0'
    },
    'bright_blue_red': {
        'primary': '#101010', 'accent': '#0354D7', 'bg': '#F2F5F6',
        'text': '#101010', 'text_muted': '#666666', 'card_bg': '#FFFFFF', 'card_border': '#E0E0E0'
    },
    'deep_red_blue': {
        'primary': '#101010', 'accent': '#731817', 'bg': '#F2F5F6',
        'text': '#101010', 'text_muted': '#666666', 'card_bg': '#FFFFFF', 'card_border': '#E0E0E0'
    },
    'deep_green_gold': {
        'primary': '#101010', 'accent': '#204023', 'bg': '#F2F5F6',
        'text': '#101010', 'text_muted': '#666666', 'card_bg': '#FFFFFF', 'card_border': '#E0E0E0'
    },
    'deep_blue_gold': {
        'primary': '#101010', 'accent': '#132D86', 'bg': '#F2F5F6',
        'text': '#101010', 'text_muted': '#666666', 'card_bg': '#FFFFFF', 'card_border': '#E0E0E0'
    },
    'blue_green_gold': {
        'primary': '#101010', 'accent': '#122C86', 'bg': '#F2F5F6',
        'text': '#101010', 'text_muted': '#666666', 'card_bg': '#FFFFFF', 'card_border': '#E0E0E0'
    },
}

# SmartArt type mappings
SMARTART_TYPE_MAP = {
    'stairs-blocks': SMARTART_TYPE.STEP_DOWN_PROCESS,
    'stairs-cubes': SMARTART_TYPE.BASIC_PROCESS,
    'stairs-focused': SMARTART_TYPE.BASIC_CHEVRON_PROCESS,
    'journey-cards': SMARTART_TYPE.BASIC_PROCESS,
    'journey-rocks': SMARTART_TYPE.ALTERNATING_FLOW,
    'funnel-diagram': SMARTART_TYPE.FUNNEL,
    'block-hierarchy': SMARTART_TYPE.HIERARCHY,
    'gem-pyramid': SMARTART_TYPE.BASIC_PYRAMID,
    'pyramid-isometric': SMARTART_TYPE.BASIC_PYRAMID,
    'pyramid-alternate-labels': SMARTART_TYPE.SEGMENTED_PYRAMID,
    'pyramid-cubes-with-arrow': SMARTART_TYPE.PYRAMID_LIST,
    'comparison-house-foundation': SMARTART_TYPE.BALANCE,
    'comparison-overlapping-cards': SMARTART_TYPE.BASIC_VENN,
    'comparison-podium-trophies': SMARTART_TYPE.BASIC_BLOCK_LIST,
    'matrix-grid-2x2': SMARTART_TYPE.BASIC_MATRIX,
    'matrix-grid-2x2-cards': SMARTART_TYPE.TITLED_MATRIX,
    'matrix-curved-quadrant': SMARTART_TYPE.BASIC_MATRIX,
    'priority-matrix': SMARTART_TYPE.TITLED_MATRIX,
    'bullseye-progression': SMARTART_TYPE.BASIC_TARGET,
    'bullseye-with-support': SMARTART_TYPE.NESTED_TARGET,
    'bullseye-single-arrow': SMARTART_TYPE.BASIC_TARGET,
    'bullseye-multi-arrows': SMARTART_TYPE.TARGET_LIST,
    'edge-analysis': SMARTART_TYPE.RADIAL,
    'edge-circular-petals': SMARTART_TYPE.RADIAL,
    'edge-hexagon-nodes': SMARTART_TYPE.RADIAL,
    'edge-rectangular-boxes': SMARTART_TYPE.BASIC_BLOCK_LIST,
    'iceberg-depth': SMARTART_TYPE.INVERTED_PYRAMID,
    'iceberg-simple-mountain': SMARTART_TYPE.BASIC_PYRAMID,
    'iceberg-complex-layers': SMARTART_TYPE.SEGMENTED_PYRAMID,
    'range-spectrum': SMARTART_TYPE.BASIC_PROCESS,
    'spectrum-vertical': SMARTART_TYPE.VERTICAL_BULLET_LIST,
    'spectrum-horizontal': SMARTART_TYPE.HORIZONTAL_BULLET_LIST,
    'brain-mapping': SMARTART_TYPE.RADIAL,
    'decision-branching': SMARTART_TYPE.HIERARCHY,
    'fishbone-diagram': SMARTART_TYPE.HIERARCHY,
    'metrics-grid': SMARTART_TYPE.BASIC_BLOCK_LIST,
    'distribution-donut': SMARTART_TYPE.BASIC_CYCLE,
}

COLOR_SCHEME_MAP = {
    'soft_peach_cream': SMARTART_COLORS.ACCENT_2,
    'executive': SMARTART_COLORS.ACCENT_1,
    'forest_green': SMARTART_COLORS.ACCENT_3,
    'sunset_orange': SMARTART_COLORS.ACCENT_2,
    'cosmic': SMARTART_COLORS.ACCENT_6,
    'azure': SMARTART_COLORS.ACCENT_1,
    'bright_red_blue': SMARTART_COLORS.COLORFUL_ACCENT_COLORS,
    'bright_blue_red': SMARTART_COLORS.COLORFUL_ACCENT_2_TO_3,
    'deep_red_blue': SMARTART_COLORS.ACCENT_6,
    'deep_green_gold': SMARTART_COLORS.ACCENT_3,
    'deep_blue_gold': SMARTART_COLORS.ACCENT_1,
    'blue_green_gold': SMARTART_COLORS.COLORFUL_ACCENT_3_TO_4,
}

SAMPLE_DATA = {
    'sequential': [
        {'label': '第一阶段', 'description': '规划与准备'},
        {'label': '第二阶段', 'description': '执行与开发'},
        {'label': '第三阶段', 'description': '测试与优化'},
        {'label': '第四阶段', 'description': '发布与维护'},
    ],
    'funnel': [
        {'label': '曝光', 'value': 10000},
        {'label': '点击', 'value': 3000},
        {'label': '注册', 'value': 800},
        {'label': '付费', 'value': 200},
    ],
    'journey': [
        {'label': '需求分析', 'description': '明确项目目标'},
        {'label': '方案设计', 'description': '制定实施计划'},
        {'label': '开发实现', 'description': '编码与集成'},
        {'label': '上线运营', 'description': '部署与监控'},
    ],
    'hierarchy': [
        {'label': '战略层', 'description': '企业愿景与目标'},
        {'label': '管理层', 'description': '资源调配与决策'},
        {'label': '执行层', 'description': '任务执行与反馈'},
        {'label': '支撑层', 'description': '基础设施与工具'},
    ],
    'comparison': [
        {'label': '方案 A', 'description': '成本低，周期长'},
        {'label': '方案 B', 'description': '成本中，周期中'},
        {'label': '方案 C', 'description': '成本高，周期短'},
    ],
}


# =============================================================================
# Slide Master Setup (L1 Layer)
# =============================================================================

def setup_slide_master(prs, config):
    """Setup slide master with theme colors, background, and decorative shapes.

    All slides will inherit these elements automatically.
    """
    theme_id = config.get('theme', 'soft_peach_cream')
    theme = THEME_COLORS.get(theme_id, THEME_COLORS['soft_peach_cream'])
    master_config = config.get('slideMaster', {})

    master = prs.slide_masters[0]

    # 1. Set theme colors
    color_scheme = master.theme.color_scheme
    color_scheme.accent1 = hex_to_rgb(theme['accent'])
    color_scheme.dk1 = hex_to_rgb(theme['primary'])
    color_scheme.lt1 = hex_to_rgb(theme['bg'])

    # 2. Set master background
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = hex_to_rgb(theme['bg'])

    # 3. Add decorative shapes to master (inherited by all slides)
    decorative_shapes = master_config.get('decorativeShapes', [])
    for shape_config in decorative_shapes:
        add_master_decorative_shape(master, shape_config, theme)

    # 4. Add placeholders (page number, date, logo) to master
    placeholders = master_config.get('placeholders', [])
    add_master_placeholders(master, placeholders, theme)

    print(f"Slide master configured: theme={theme_id}, shapes={len(decorative_shapes)}, placeholders={placeholders}")


def add_master_decorative_shape(master, shape_config: dict, theme: dict):
    """Add decorative shape to slide master."""
    shape_id = shape_config.get('id')
    preset = shape_config.get('preset', 'default')
    accent_color = hex_to_rgb(theme['accent'])

    if shape_id == 'corner-accent':
        # Top-right corner accent badge
        if preset in ['top-right', 'default']:
            badge = master.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.2), Inches(0.3), Inches(1.5), Inches(0.5)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = accent_color
            badge.line.fill.background()

            # Decorative dots below badge
            for i in range(6):
                dot = master.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(8.4 + i * 0.15), Inches(0.9), Inches(0.06), Inches(0.06)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = accent_color
                dot.line.fill.background()

    elif shape_id == 'side-strip':
        # Left side accent strip
        if preset in ['left', 'default']:
            strip = master.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(0.08), Inches(7.5)
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = accent_color
            strip.line.fill.background()

    elif shape_id == 'bottom-line':
        # Bottom decorative line
        line = master.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4), Inches(7.2), Inches(9.2), Inches(0.015)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(theme.get('card_border', '#E0E0E0'))
        line.line.fill.background()


def add_master_placeholders(master, placeholders: list, theme: dict):
    """Add placeholder shapes to slide master.

    TODO: Currently uses textboxes as workaround.
    Should use master.shapes.add_placeholder() when python-pptx supports it.
    """
    muted_color = hex_to_rgb(theme.get('text_muted', '#888888'))

    if 'page-number' in placeholders:
        # TODO: Should be master.shapes.add_placeholder(PP_PLACEHOLDER.SLIDE_NUMBER, ...)
        ph = master.shapes.add_textbox(
            Inches(9), Inches(7.1), Inches(0.8), Inches(0.3)
        )
        p = ph.text_frame.paragraphs[0]
        p.text = "<#>"
        p.font.size = Pt(10)
        p.font.color.rgb = muted_color
        p.alignment = PP_ALIGN.RIGHT

    if 'date' in placeholders:
        # TODO: Should be master.shapes.add_placeholder(PP_PLACEHOLDER.DATE, ...)
        ph = master.shapes.add_textbox(
            Inches(0.4), Inches(7.1), Inches(2), Inches(0.3)
        )
        p = ph.text_frame.paragraphs[0]
        p.text = "2024年12月"
        p.font.size = Pt(10)
        p.font.color.rgb = muted_color

    if 'logo' in placeholders:
        # TODO: Should be master.shapes.add_placeholder(PP_PLACEHOLDER.PICTURE, ...)
        ph = master.shapes.add_textbox(
            Inches(0.4), Inches(0.2), Inches(1.5), Inches(0.3)
        )
        p = ph.text_frame.paragraphs[0]
        p.text = "LOGO"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = muted_color


# =============================================================================
# Slide Content Helpers
# =============================================================================

def add_title_with_tag(slide, title: str, tag: str, theme: dict, title_style: str = 'with-tag'):
    """Add title with optional tag badge."""
    accent_color = hex_to_rgb(theme['accent'])
    text_color = hex_to_rgb(theme['primary'])

    if title_style == 'with-tag' and tag:
        # Tag badge
        tag_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.4), Inches(0.3), Inches(0.8), Inches(0.35)
        )
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = accent_color
        tag_shape.line.fill.background()

        tag_tf = tag_shape.text_frame
        tag_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tag_p = tag_tf.paragraphs[0]
        tag_p.text = tag
        tag_p.font.size = Pt(10)
        tag_p.font.bold = True
        tag_p.font.color.rgb = RGBColor(255, 255, 255)
        tag_tf.margin_top = Pt(4)
        tag_tf.margin_bottom = Pt(4)

        title_y = Inches(0.75)
    else:
        title_y = Inches(0.4)

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.4), title_y, Inches(7), Inches(0.7)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = text_color


def add_source_citation(slide, source: str, theme: dict):
    """Add source citation at bottom."""
    text_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(7.1), Inches(5), Inches(0.3)
    )
    p = text_box.text_frame.paragraphs[0]
    p.text = f"数据来源：{source}"
    p.font.size = Pt(9)
    p.font.color.rgb = hex_to_rgb(theme.get('text_muted', '#888888'))


def create_smartart_data(data_type: str, smartart_type_id: str) -> SmartArtData:
    """Create SmartArtData based on data type."""
    data = SmartArtData()
    samples = SAMPLE_DATA.get(data_type, SAMPLE_DATA['comparison'])

    pptx_type = SMARTART_TYPE_MAP.get(smartart_type_id)
    is_hierarchy = pptx_type in [SMARTART_TYPE.HIERARCHY, SMARTART_TYPE.ORGANIZATION_CHART]
    is_radial = pptx_type in [SMARTART_TYPE.RADIAL]

    if is_hierarchy and len(samples) > 1:
        root = data.add_node(samples[0]['label'])
        for item in samples[1:]:
            root.add_child(item['label'])
    elif is_radial and len(samples) > 1:
        center = data.add_node(samples[0]['label'])
        for item in samples[1:]:
            center.add_child(item['label'])
    else:
        for item in samples:
            text = item['label']
            if item.get('description'):
                text += f"\n{item['description']}"
            data.add_node(text)

    return data


# =============================================================================
# Slide Generators
# =============================================================================

def generate_cover_slide(prs, config, theme):
    """Generate cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    accent_color = hex_to_rgb(theme['accent'])
    text_color = hex_to_rgb(theme['primary'])

    # Decorative accent shape (top right)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(0.5), Inches(2), Inches(3)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5), Inches(6.5), Inches(1.5)
    )
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.text = "演示文稿标题"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = text_color

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.2), Inches(6.5), Inches(0.8)
    )
    p = sub_box.text_frame.paragraphs[0]
    p.text = "副标题或作者信息"
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(theme.get('text_muted', '#888888'))

    # Date
    date_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.5), Inches(6), Inches(0.4)
    )
    p = date_box.text_frame.paragraphs[0]
    p.text = "2024年12月"
    p.font.size = Pt(12)
    p.font.color.rgb = hex_to_rgb(theme.get('text_muted', '#888888'))

    return slide


def generate_divider_slide(prs, config, theme):
    """Generate divider/section slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    divider_config = config.get('divider', {})
    section_index = divider_config.get('sectionIndex', 1)
    layout_type = divider_config.get('layout', 'left-align')

    sections = [
        {'num': '01', 'zh': '年度工作概述', 'en': 'Annual Overview', 'part': '第一部分'},
        {'num': '02', 'zh': '工作完成情况', 'en': 'Work Progress', 'part': '第二部分'},
        {'num': '03', 'zh': '项目成果展示', 'en': 'Project Results', 'part': '第三部分'},
        {'num': '04', 'zh': '工作不足与改进', 'en': 'Improvements', 'part': '第四部分'},
    ]

    idx = section_index - 1 if section_index > 0 else 0
    section = sections[idx % len(sections)]
    accent_color = hex_to_rgb(theme['accent'])

    # Set background based on layout type
    if layout_type in ['left-align', 'fullbleed']:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = accent_color
        text_color = RGBColor(255, 255, 255)
        strip_color = RGBColor(255, 255, 255)
    else:
        text_color = hex_to_rgb(theme['primary'])
        strip_color = accent_color

    # Left strip
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.1), Inches(7.5)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = strip_color
    strip.line.fill.background()

    # Part label
    part_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(3), Inches(0.8))
    p = part_box.text_frame.paragraphs[0]
    p.text = section['part']
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(7), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = section['zh']
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = text_color

    # English subtitle
    en_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(5), Inches(0.5))
    p = en_box.text_frame.paragraphs[0]
    p.text = section['en']
    p.font.size = Pt(14)
    p.font.color.rgb = text_color

    # Large number
    num_box = slide.shapes.add_textbox(Inches(6), Inches(1.5), Inches(3.5), Inches(5))
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    p = num_box.text_frame.paragraphs[0]
    p.text = section['num']
    p.font.size = Pt(280)
    p.font.color.rgb = text_color if layout_type in ['left-align', 'fullbleed'] else accent_color

    return slide


def generate_smartart_slide(prs, config, theme):
    """Generate SmartArt content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    smartart_config = config.get('smartart', {})
    smartart_type_id = smartart_config.get('type', 'funnel-diagram')
    data_type = smartart_config.get('dataType', 'comparison')
    placement = smartart_config.get('placement', 'full')
    theme_id = config.get('theme', 'azure')

    # Title
    add_title_with_tag(slide, "流程分析", "SmartArt", theme, 'with-tag')
    add_source_citation(slide, "内部分析报告", theme)

    # SmartArt area
    content_top = Inches(1.4)
    content_height = Inches(5.5)
    text_color = hex_to_rgb(theme['text'])

    if placement == 'full':
        x, y, cx, cy = Inches(0.4), content_top, Inches(9.2), content_height
    elif placement == 'left-desc':
        x, y, cx, cy = Inches(0.4), content_top, Inches(5.5), content_height
        # Description card
        desc_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.1), content_top, Inches(3.5), content_height
        )
        desc_card.fill.solid()
        desc_card.fill.fore_color.rgb = hex_to_rgb(theme.get('card_bg', '#FFFFFF'))
        desc_card.line.color.rgb = hex_to_rgb(theme.get('card_border', '#E0E0E0'))

        header = slide.shapes.add_textbox(Inches(6.3), content_top + Inches(0.2), Inches(3.1), Inches(0.4))
        p = header.text_frame.paragraphs[0]
        p.text = "流程说明"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color

        desc = slide.shapes.add_textbox(Inches(6.3), content_top + Inches(0.7), Inches(3.1), Inches(4))
        desc.text_frame.word_wrap = True
        p = desc.text_frame.paragraphs[0]
        p.text = "这里可以放置对 SmartArt 图形的描述。"
        p.font.size = Pt(11)
        p.font.color.rgb = text_color
    else:
        x, y, cx, cy = Inches(0.4), content_top, Inches(9.2), content_height

    # Add SmartArt
    pptx_type = SMARTART_TYPE_MAP.get(smartart_type_id, SMARTART_TYPE.BASIC_PROCESS)
    color_scheme = COLOR_SCHEME_MAP.get(theme_id, SMARTART_COLORS.ACCENT_1)
    smartart_data = create_smartart_data(data_type, smartart_type_id)

    try:
        slide.shapes.add_smartart(pptx_type, x, y, cx, cy, smartart_data, color_scheme=color_scheme)
        print(f"Added SmartArt: {smartart_type_id}")
    except Exception as e:
        print(f"SmartArt failed: {e}")
        # Fallback placeholder
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(theme.get('card_bg', '#FFFFFF'))
        card.line.color.rgb = hex_to_rgb(theme.get('card_border', '#E0E0E0'))

    return slide


def generate_grid_slide(prs, config, theme):
    """Generate grid/content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    layout_config = config.get('layout', {})
    layout_type = layout_config.get('type', 'chart-left')
    title_style = layout_config.get('titleStyle', 'with-tag')
    source_style = layout_config.get('sourceStyle', 'citation')

    # Title
    add_title_with_tag(slide, "市场趋势分析", "分析报告", theme, title_style)

    if source_style == 'citation':
        add_source_citation(slide, "行业研究报告 2024", theme)

    content_top = Inches(1.4) if title_style == 'with-tag' else Inches(1.1)
    content_height = Inches(5.5)

    accent_color = hex_to_rgb(theme['accent'])
    text_color = hex_to_rgb(theme['text'])
    card_bg = hex_to_rgb(theme.get('card_bg', '#FFFFFF'))
    card_border = hex_to_rgb(theme.get('card_border', '#E0E0E0'))

    if 'left' in layout_type:
        # Left card
        left = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.4), content_top, Inches(4.4), content_height
        )
        left.fill.solid()
        left.fill.fore_color.rgb = card_bg
        left.line.color.rgb = card_border

        header = slide.shapes.add_textbox(Inches(0.6), content_top + Inches(0.2), Inches(4), Inches(0.4))
        p = header.text_frame.paragraphs[0]
        p.text = "关键指标"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color

        # Right card
        right = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5), content_top, Inches(4.4), content_height
        )
        right.fill.solid()
        right.fill.fore_color.rgb = card_bg
        right.line.color.rgb = card_border

        header = slide.shapes.add_textbox(Inches(5.2), content_top + Inches(0.2), Inches(4), Inches(0.4))
        p = header.text_frame.paragraphs[0]
        p.text = "市场分析"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color

        # Bullet points
        bullets = [
            "市场规模持续增长，年复合增长率达 15%",
            "用户活跃度显著提升，日活突破百万",
            "产品线扩展计划已完成第一阶段"
        ]
        for i, bullet in enumerate(bullets):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(5.3), content_top + Inches(0.85 + i * 0.6), Inches(0.12), Inches(0.12)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent_color
            dot.line.fill.background()

            text = slide.shapes.add_textbox(
                Inches(5.55), content_top + Inches(0.75 + i * 0.6), Inches(3.8), Inches(0.5)
            )
            text.text_frame.word_wrap = True
            p = text.text_frame.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(12)
            p.font.color.rgb = text_color
    else:
        # Single card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.4), content_top, Inches(9.2), content_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = card_border

    return slide


# =============================================================================
# Main Generator
# =============================================================================

def generate_pptx(config: dict, output_path: str):
    """Generate PPTX from config using Slide Master API."""
    prs = Presentation.create()

    theme_id = config.get('theme', 'soft_peach_cream')
    theme = THEME_COLORS.get(theme_id, THEME_COLORS['soft_peach_cream'])

    # Setup slide master first (L1 layer)
    setup_slide_master(prs, config)

    # Generate slide based on page type
    page_type = config.get('pageType', 'content-grid')

    if page_type == 'cover':
        generate_cover_slide(prs, config, theme)
    elif page_type == 'divider':
        generate_divider_slide(prs, config, theme)
    elif page_type == 'content-smartart':
        generate_smartart_slide(prs, config, theme)
    else:
        generate_grid_slide(prs, config, theme)

    prs.save(output_path)
    print(f"Saved: {output_path}")


def main():
    if len(sys.argv) < 2:
        demo_config = {
            "theme": "soft_peach_cream",
            "pageType": "content-grid",
            "slideMaster": {
                "decorativeShapes": [
                    {"id": "corner-accent", "preset": "top-right"},
                    {"id": "side-strip", "preset": "left"}
                ],
                "placeholders": ["page-number"]
            },
            "layout": {
                "type": "chart-left",
                "titleStyle": "with-tag",
                "sourceStyle": "citation"
            }
        }
        output_path = "/Users/touichskyer/Desktop/layout-template-generator/output.pptx"
        generate_pptx(demo_config, output_path)
        print(f"Demo generated: {output_path}")
        return

    config_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"

    with open(config_path, 'r', encoding='utf-8') as f:
        config = json.load(f)

    generate_pptx(config, output_path)


if __name__ == '__main__':
    main()
