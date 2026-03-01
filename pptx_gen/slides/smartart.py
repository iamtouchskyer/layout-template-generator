"""
SmartArt content slide generator.
"""

import logging

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.smartart import SMARTART_TYPE, SMARTART_COLORS
from pptx.smartart import SmartArtData

from ..themes import hex_to_rgb, SMARTART_TYPE_MAP, SMARTART_COLOR_SCHEME_MAP
from ..dimensions import px_to_inches_x, px_to_inches_y
from ..smartart_layout_mapper import (
    resolve_ambiguous_type_ids_from_pptx_enum,
    resolve_type_id_from_layout_id,
    resolve_type_id_from_pptx_enum,
)
from ..utils import get_blank_layout, remove_slide_placeholders
from .grid import (
    add_source_citation_dynamic,
    resolve_content_shell_geometry,
    set_title_with_style,
)


def generate_smartart_slide(prs, config, theme):
    """Generate SmartArt content slide."""
    slide = prs.slides.add_slide(get_blank_layout(prs))
    remove_slide_placeholders(slide)

    smartart_config = config.get('smartart', {})
    smartart_type_id = _resolve_smartart_type_id(smartart_config)
    placement = smartart_config.get('placement', 'full')
    color_scheme_id = smartart_config.get('colorScheme', 'colorful2')

    items = _extract_smartart_items(smartart_config)

    shell = resolve_content_shell_geometry(config)
    title_style = shell["title_style"]
    has_title = shell["has_title"]
    has_source = shell["has_source"]
    header_x = shell["header_in"]["x"]
    header_y = shell["header_in"]["y"]
    header_w = shell["header_in"]["w"]
    header_h = shell["header_in"]["h"]
    body_x = shell["body_in"]["x"]
    body_y = shell["body_in"]["y"]
    body_w = shell["body_in"]["w"]
    body_h = shell["body_in"]["h"]
    footer_x = shell["footer_in"]["x"]
    footer_y = shell["footer_in"]["y"]
    footer_w = shell["footer_in"]["w"]
    content_title = shell["content_title"]
    content_tag = shell["content_tag"]
    content_source = shell["content_source"]

    if has_title:
        set_title_with_style(
            slide,
            None,
            content_title,
            content_tag,
            theme,
            title_style,
            bounds=(header_x, header_y, header_w, header_h),
        )

    if has_source:
        add_source_citation_dynamic(
            slide,
            content_source,
            theme,
            x=footer_x,
            y=footer_y,
            width=footer_w,
        )

    # SmartArt area (inside body bounds)
    gap = px_to_inches_x(32)
    desc_width = max(px_to_inches_x(260), min(px_to_inches_x(360), body_w * 0.34))
    desc_height = max(px_to_inches_y(72), min(px_to_inches_y(100), body_h * 0.24))
    min_chart_w = px_to_inches_x(180)
    min_chart_h = px_to_inches_y(160)
    text_color = hex_to_rgb(theme['text'])

    if placement == 'full':
        x, y, cx, cy = Inches(body_x), Inches(body_y), Inches(body_w), Inches(body_h)
    elif placement == 'left-desc':
        max_desc_w = max(px_to_inches_x(80), body_w - gap - min_chart_w)
        desc_w = min(desc_width, max_desc_w)
        chart_w = max(min_chart_w, body_w - desc_w - gap)
        desc_x = body_x + chart_w + gap
        x, y, cx, cy = Inches(body_x), Inches(body_y), Inches(chart_w), Inches(body_h)
        _add_description_block(slide, Inches(desc_x), Inches(body_y), Inches(desc_w), Inches(body_h), text_color)
    elif placement == 'right-desc':
        max_desc_w = max(px_to_inches_x(80), body_w - gap - min_chart_w)
        desc_w = min(desc_width, max_desc_w)
        chart_w = max(min_chart_w, body_w - desc_w - gap)
        chart_x = body_x + desc_w + gap
        x, y, cx, cy = Inches(chart_x), Inches(body_y), Inches(chart_w), Inches(body_h)
        _add_description_block(slide, Inches(body_x), Inches(body_y), Inches(desc_w), Inches(body_h), text_color)
    elif placement == 'top-desc':
        max_desc_h = max(px_to_inches_y(40), body_h - gap - min_chart_h)
        desc_h = min(desc_height, max_desc_h)
        chart_h = max(min_chart_h, body_h - desc_h - gap)
        chart_y = body_y + desc_h + gap
        x, y, cx, cy = Inches(body_x), Inches(chart_y), Inches(body_w), Inches(chart_h)
        _add_description_block(slide, Inches(body_x), Inches(body_y), Inches(body_w), Inches(desc_h), text_color, compact=True)
    elif placement == 'bottom-desc':
        max_desc_h = max(px_to_inches_y(40), body_h - gap - min_chart_h)
        desc_h = min(desc_height, max_desc_h)
        chart_h = max(min_chart_h, body_h - desc_h - gap)
        desc_y = body_y + chart_h + gap
        x, y, cx, cy = Inches(body_x), Inches(body_y), Inches(body_w), Inches(chart_h)
        _add_description_block(slide, Inches(body_x), Inches(desc_y), Inches(body_w), Inches(desc_h), text_color, compact=True)
    else:
        x, y, cx, cy = Inches(body_x), Inches(body_y), Inches(body_w), Inches(body_h)

    # Add SmartArt with selected color scheme
    pptx_type = SMARTART_TYPE_MAP.get(smartart_type_id, SMARTART_TYPE.BASIC_PYRAMID)
    color_scheme = SMARTART_COLOR_SCHEME_MAP.get(color_scheme_id, SMARTART_COLORS.COLORFUL_ACCENT_COLORS)
    smartart_data = _create_smartart_data(smartart_type_id, items)

    try:
        slide.shapes.add_smartart(pptx_type, x, y, cx, cy, smartart_data, color_scheme=color_scheme)
    except (KeyError, ValueError, AttributeError, TypeError) as e:
        logging.warning(f"SmartArt creation failed for type '{smartart_type_id}': {e}")
        _add_fallback_placeholder(slide, x, y, cx, cy, theme)

    return slide


def _resolve_smartart_type_id(smartart_config: dict) -> str:
    """Resolve SmartArt type id with backward-compatible fallbacks.

    Priority:
    1. smartart.ooxml.layoutId
    2. smartart.typeId (preferred external API)
    3. smartart.type when it is internal type-id
    4. smartart.type when it is SMARTART_TYPE enum name
    """
    if isinstance(smartart_config, dict):
        ooxml_data = smartart_config.get('ooxml')
        if isinstance(ooxml_data, dict):
            layout_id = ooxml_data.get('layoutId')
            mapped = resolve_type_id_from_layout_id(layout_id)
            if mapped:
                return mapped

        explicit_type_id = smartart_config.get('typeId')
        if isinstance(explicit_type_id, str):
            normalized_type_id = explicit_type_id.strip()
            if normalized_type_id in SMARTART_TYPE_MAP:
                return normalized_type_id
            if normalized_type_id:
                logging.warning(
                    "Unknown SmartArt typeId '%s'; falling back to smartart.type.",
                    normalized_type_id,
                )

        explicit_type = smartart_config.get('type')
        if isinstance(explicit_type, str):
            normalized_type = explicit_type.strip()

            if normalized_type in SMARTART_TYPE_MAP:
                return normalized_type

            mapped_from_enum = resolve_type_id_from_pptx_enum(normalized_type)
            if mapped_from_enum:
                ambiguous_candidates = resolve_ambiguous_type_ids_from_pptx_enum(normalized_type)
                if ambiguous_candidates:
                    logging.warning(
                        "Ambiguous SMARTART_TYPE '%s' resolved to '%s'. Candidates: %s. "
                        "Prefer smartart.typeId or smartart.ooxml.layoutId to disambiguate.",
                        normalized_type,
                        mapped_from_enum,
                        ", ".join(ambiguous_candidates),
                    )
                return mapped_from_enum

            if normalized_type:
                logging.warning("Unknown SmartArt type '%s'; falling back to default.", normalized_type)

    return 'pyramid'


def _add_title_with_tag(slide, title: str, tag: str, theme: dict, title_style: str):
    """Add title with optional tag badge."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    accent_color = hex_to_rgb(theme['accent'])
    text_color = hex_to_rgb(theme['primary'])

    if title_style == 'with-tag' and tag:
        tag_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.4), Inches(0.3), Inches(0.8), Inches(0.35)
        )
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = accent_color
        tag_shape.line.fill.background()

        tag_tf = tag_shape.text_frame
        tag_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tag_p = tag_tf.paragraphs[0]
        tag_p.text = tag
        tag_p.font.size = Pt(10)
        tag_p.font.bold = True
        tag_p.font.color.rgb = RGBColor(255, 255, 255)
        tag_tf.margin_top = Pt(4)
        tag_tf.margin_bottom = Pt(4)

        title_y = Inches(0.75)
    else:
        title_y = Inches(0.4)

    title_box = slide.shapes.add_textbox(
        Inches(0.4), title_y, Inches(7), Inches(0.7)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = text_color


def _add_source_citation(slide, source: str, theme: dict):
    """Add source citation at bottom."""
    text_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(7.1), Inches(5), Inches(0.3)
    )
    p = text_box.text_frame.paragraphs[0]
    p.text = f"数据来源：{source}"
    p.font.size = Pt(9)
    p.font.color.rgb = hex_to_rgb(theme.get('text_muted', '#888888'))


def _add_description_block(slide, x, y, w, h, text_color, compact=False):
    """Add description text block without background card (match preview style)."""
    header = slide.shapes.add_textbox(x, y + Inches(0.1), w, Inches(0.4))
    p = header.text_frame.paragraphs[0]
    p.text = "说明文字区域"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = text_color

    desc = slide.shapes.add_textbox(x, y + Inches(0.55), w, h - Inches(0.55))
    desc.text_frame.word_wrap = True
    p = desc.text_frame.paragraphs[0]
    p.text = "这里可以放置对 SmartArt 图形的描述、解释或相关数据说明。"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    if compact:
        return

    for line in ["要点一：关键信息", "要点二：补充说明", "要点三：总结概括"]:
        pp = desc.text_frame.add_paragraph()
        pp.text = f"• {line}"
        pp.font.size = Pt(11)
        pp.font.color.rgb = text_color


def _add_fallback_placeholder(slide, x, y, cx, cy, theme):
    """Add fallback placeholder when SmartArt fails."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(theme.get('card_bg', '#FFFFFF'))
    card.line.color.rgb = hex_to_rgb(theme.get('card_border', '#E0E0E0'))


def _create_smartart_data(smartart_type_id: str, items: list) -> SmartArtData:
    """Create SmartArtData from items list."""
    data = SmartArtData()

    if not items:
        items = ['项目 1', '项目 2', '项目 3', '项目 4']

    pptx_type = SMARTART_TYPE_MAP.get(smartart_type_id)
    is_hierarchy = pptx_type in [SMARTART_TYPE.HIERARCHY, SMARTART_TYPE.ORGANIZATION_CHART]
    is_radial = pptx_type in [SMARTART_TYPE.BASIC_RADIAL, SMARTART_TYPE.RADIAL]
    supports_nested_children = smartart_type_id in {'pyramid', 'pyramid-list', 'cycle'}

    if (is_hierarchy or is_radial or supports_nested_children) and _has_nested_children(items):
        for item in items:
            root = data.add_node(_item_text(item))
            if isinstance(item, dict):
                _add_child_nodes(root, item.get('children', []))
    elif is_hierarchy and len(items) > 1:
        root = data.add_node(_item_text(items[0]))
        for item in items[1:]:
            root.add_child(_item_text(item))
    elif is_radial and len(items) > 1:
        center = data.add_node(_item_text(items[0]))
        for item in items[1:]:
            center.add_child(_item_text(item))
    else:
        for item in items:
            data.add_node(_item_text(item))

    return data


def _extract_smartart_items(smartart_config: dict) -> list:
    """Extract SmartArt items from config with explicit items first, OOXML fallback."""
    items = smartart_config.get('items')
    if not isinstance(items, list) or len(items) == 0:
        ooxml_data = smartart_config.get('ooxml', {})
        if isinstance(ooxml_data, dict):
            items = ooxml_data.get('items', [])
        else:
            items = []
    return items if isinstance(items, list) else []


def _item_text(item) -> str:
    """Return display text from SmartArt item."""
    if isinstance(item, dict):
        return str(item.get('text', str(item)))
    return str(item)


def _has_nested_children(items: list) -> bool:
    """Check whether items include explicit nested children."""
    for item in items:
        if isinstance(item, dict):
            children = item.get('children', [])
            if isinstance(children, list) and len(children) > 0:
                return True
    return False


def _add_child_nodes(parent_node, children: list) -> None:
    """Recursively add child nodes from nested items."""
    if not isinstance(children, list):
        return
    for child in children:
        node = parent_node.add_child(_item_text(child))
        if isinstance(child, dict):
            _add_child_nodes(node, child.get('children', []))
