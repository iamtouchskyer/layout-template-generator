"""
Slide Master configuration for PPTX generation.
"""

from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER

from .themes import hex_to_rgb
from .dimensions import set_slide_dimensions, px_to_inches_x, px_to_inches_y, get_pptx_dimensions
from .decorative_shapes import add_master_decorative_shape
from .master_placeholders import add_master_placeholders

# Re-export for backward compatibility
__all__ = ['setup_slide_master', 'set_slide_dimensions', 'px_to_inches_x', 'px_to_inches_y', 'get_pptx_dimensions']


def setup_slide_master(prs, config):
    """Setup slide master with background, placeholders, and decorative shapes.

    Theme colors are already set via Presentation.create(theme=...).
    This function handles: slide dimensions, background, placeholder positions,
    decorative shapes, and master placeholders.
    """
    from .themes import get_theme

    # 0. Set slide dimensions from config (for coordinate conversion)
    slide_config = config.get('slide', {})
    if slide_config:
        frontend_w = slide_config.get('width', 1280)
        frontend_h = slide_config.get('height', 720)
        pptx_w = slide_config.get('widthInches', 13.333)
        pptx_h = slide_config.get('heightInches', 7.5)
        set_slide_dimensions(frontend_w, frontend_h, pptx_w, pptx_h)

        # Set presentation slide size
        prs.slide_width = Inches(pptx_w)
        prs.slide_height = Inches(pptx_h)

    theme_id = config.get('theme', 'soft_peach_cream')
    theme = get_theme(theme_id)
    master_config = config.get('slideMaster', {})
    content_areas = master_config.get('contentAreas', {})

    master = prs.slide_masters[0]

    # 1. Set master background
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = hex_to_rgb(theme['bg'])

    # 3. Use pre-calculated bounds from frontend (in pixels)
    # Coordinate conversion uses dimensions set above
    header_bounds = content_areas.get('headerBounds', {})
    body_bounds = content_areas.get('bodyBounds', {})
    title_style = content_areas.get('titleStyle', 'with-tag')

    # 4. Adjust title placeholder position using pre-calculated bounds
    title_ph = master.placeholders.get(PP_PLACEHOLDER.TITLE)
    if title_ph and header_bounds:
        title_ph.left = Inches(px_to_inches_x(header_bounds.get('x', 40)))
        title_ph.top = Inches(px_to_inches_y(header_bounds.get('y', 20)))
        title_ph.width = Inches(px_to_inches_x(header_bounds.get('width', 1200)))
        title_ph.height = Inches(px_to_inches_y(header_bounds.get('height', 60)))

    # 5. Adjust body placeholder position using pre-calculated bounds
    body_ph = master.placeholders.get(PP_PLACEHOLDER.BODY)
    if body_ph and body_bounds:
        body_ph.left = Inches(px_to_inches_x(body_bounds.get('x', 40)))
        body_ph.top = Inches(px_to_inches_y(body_bounds.get('y', 88)))
        body_ph.width = Inches(px_to_inches_x(body_bounds.get('width', 1200)))
        body_ph.height = Inches(px_to_inches_y(body_bounds.get('height', 592)))

    # 6. Add decorative shapes to master (inherited by all slides)
    decorative_shapes = master_config.get('decorativeShapes', [])
    for shape_config in decorative_shapes:
        add_master_decorative_shape(master, shape_config, theme)

    # 7. Add placeholders (page number, date, logo) to master
    placeholders = master_config.get('placeholders', {})
    add_master_placeholders(master, placeholders, theme)

