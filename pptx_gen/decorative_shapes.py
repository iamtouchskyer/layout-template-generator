"""
Decorative shapes for slide master (side-bar, corner, dots, badge, lines).
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from .themes import hex_to_rgb
from .dimensions import px_to_inches_x, px_to_inches_y, get_pptx_dimensions


def add_master_decorative_shape(master, shape_config: dict, theme: dict):
    """Add decorative shape to slide master.

    Supports two config types:
    1. Legacy preset-based: {id, preset}
    2. New thickness-positions: {id, configType, thickness, positions, thicknessConfig}
    """
    shape_id = shape_config.get('id')
    config_type = shape_config.get('configType', 'presets')
    accent_color = hex_to_rgb(theme['accent'])

    SLIDE_WIDTH, SLIDE_HEIGHT = get_pptx_dimensions()

    if config_type == 'thickness-positions':
        _add_thickness_position_shape(master, shape_id, shape_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)
    else:
        _add_preset_shape(master, shape_id, shape_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)


def _add_thickness_position_shape(master, shape_id, shape_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add shape with thickness-positions config."""
    thickness = shape_config.get('thickness', 'thin')
    positions = shape_config.get('positions', [])
    thickness_config = shape_config.get('thicknessConfig', {})

    for pos in positions:
        if shape_id == 'side-bar':
            _add_side_bar(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)
        elif shape_id == 'corner':
            _add_corner_triangle(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)
        elif shape_id == 'corner-dots':
            _add_corner_dots(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)
        elif shape_id == 'accent-circle':
            _add_accent_circle(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)
        elif shape_id == 'accent-ring':
            _add_accent_ring(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)


def _add_side_bar(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add side bar shape."""
    size = px_to_inches_x(thickness_config.get('size', 8))
    if pos == 'left':
        shape = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(size), Inches(SLIDE_HEIGHT))
    elif pos == 'right':
        shape = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(SLIDE_WIDTH - size), Inches(0), Inches(size), Inches(SLIDE_HEIGHT))
    elif pos == 'top':
        shape = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH), Inches(size))
    elif pos == 'bottom':
        shape = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(SLIDE_HEIGHT - size), Inches(SLIDE_WIDTH), Inches(size))
    else:
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()


def _add_corner_triangle(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add corner triangle shape."""
    size = px_to_inches_x(thickness_config.get('size', 60))
    if pos == 'tl':
        shape = master.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(0), Inches(size), Inches(size))
        shape.rotation = 0
    elif pos == 'tr':
        shape = master.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(SLIDE_WIDTH - size), Inches(0), Inches(size), Inches(size))
        shape.rotation = 90
    elif pos == 'bl':
        shape = master.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(SLIDE_HEIGHT - size), Inches(size), Inches(size))
        shape.rotation = 270
    elif pos == 'br':
        shape = master.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(SLIDE_WIDTH - size), Inches(SLIDE_HEIGHT - size), Inches(size), Inches(size))
        shape.rotation = 180
    else:
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()


def _add_corner_dots(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add corner dots pattern."""
    size = px_to_inches_x(thickness_config.get('size', 40))
    dot_size = px_to_inches_x(thickness_config.get('dotSize', 2))
    gap = px_to_inches_x(thickness_config.get('gap', 8))

    if pos == 'tl':
        base_x, base_y = 0.2, 0.2
    elif pos == 'tr':
        base_x, base_y = SLIDE_WIDTH - size - 0.2, 0.2
    elif pos == 'bl':
        base_x, base_y = 0.2, SLIDE_HEIGHT - size - 0.2
    elif pos == 'br':
        base_x, base_y = SLIDE_WIDTH - size - 0.2, SLIDE_HEIGHT - size - 0.2
    else:
        return

    cols = int(size / (dot_size + gap))
    rows = int(size / (dot_size + gap))
    for row in range(min(rows, 5)):
        for col in range(min(cols, 5)):
            dot = master.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(base_x + col * (dot_size + gap)),
                Inches(base_y + row * (dot_size + gap)),
                Inches(dot_size), Inches(dot_size)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent_color
            dot.line.fill.background()


def _add_accent_circle(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add accent circle shape."""
    size = px_to_inches_x(thickness_config.get('size', 100))
    offset = px_to_inches_x(thickness_config.get('offset', -25))

    if pos == 'tl':
        x, y = offset, offset
    elif pos == 'tr':
        x, y = SLIDE_WIDTH - size + abs(offset), offset
    elif pos == 'bl':
        x, y = offset, SLIDE_HEIGHT - size + abs(offset)
    elif pos == 'br':
        x, y = SLIDE_WIDTH - size + abs(offset), SLIDE_HEIGHT - size + abs(offset)
    else:
        return

    circle = master.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()


def _add_accent_ring(master, pos, thickness_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add accent ring shape."""
    size = px_to_inches_x(thickness_config.get('size', 80))
    stroke_width = thickness_config.get('strokeWidth', 3)
    offset = px_to_inches_x(thickness_config.get('offset', -20))

    if pos == 'tl':
        x, y = offset, offset
    elif pos == 'tr':
        x, y = SLIDE_WIDTH - size + abs(offset), offset
    elif pos == 'bl':
        x, y = offset, SLIDE_HEIGHT - size + abs(offset)
    elif pos == 'br':
        x, y = SLIDE_WIDTH - size + abs(offset), SLIDE_HEIGHT - size + abs(offset)
    else:
        return

    ring = master.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    ring.fill.background()
    ring.line.color.rgb = accent_color
    ring.line.width = Pt(stroke_width)


def _add_preset_shape(master, shape_id, shape_config, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add shape with legacy preset config."""
    preset = shape_config.get('preset', 'default')
    position_config = shape_config.get('positionConfig', {})

    if shape_id == 'header-badge':
        _add_header_badge(master, preset, position_config, accent_color, SLIDE_WIDTH)
    elif shape_id == 'header-line':
        _add_header_line(master, preset, accent_color, SLIDE_WIDTH)
    elif shape_id == 'footer-line':
        _add_footer_line(master, preset, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT)


def _add_header_badge(master, preset, position_config, accent_color, SLIDE_WIDTH):
    """Add header badge shape with configurable position."""
    position_config = position_config or {}

    width_px = position_config.get('width', 140)
    height_px = position_config.get('height', 40)
    width = px_to_inches_x(width_px)
    height = px_to_inches_y(height_px)

    top_px = position_config.get('top', 20)
    y = px_to_inches_y(top_px)

    if 'left' in position_config:
        x = px_to_inches_x(position_config['left'])
    elif 'right' in position_config:
        right_px = position_config['right']
        x = SLIDE_WIDTH - width - px_to_inches_x(right_px)
    else:
        x = SLIDE_WIDTH - width - 0.4

    if preset in ['pill-right', 'pill-left']:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE

    badge = master.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent_color
    badge.line.fill.background()

    badge_text = position_config.get('text', 'SECTION')
    badge.text_frame.paragraphs[0].text = badge_text
    badge.text_frame.paragraphs[0].font.size = Pt(10)
    badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_header_line(master, preset, accent_color, SLIDE_WIDTH):
    """Add header line shape."""
    top = px_to_inches_y(82)
    if preset == 'full':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), Inches(SLIDE_WIDTH), Inches(0.02))
    elif preset == 'center':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(top), Inches(SLIDE_WIDTH - 0.8), Inches(0.02))
    elif preset == 'left-half':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), Inches(SLIDE_WIDTH / 2), Inches(0.02))
    elif preset == 'right-half':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(SLIDE_WIDTH / 2), Inches(top), Inches(SLIDE_WIDTH / 2), Inches(0.02))
    else:
        return
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()


def _add_footer_line(master, preset, accent_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add footer line shape."""
    bottom = SLIDE_HEIGHT - px_to_inches_y(62)
    if preset == 'full':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(bottom), Inches(SLIDE_WIDTH), Inches(0.02))
    elif preset == 'center':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(bottom), Inches(SLIDE_WIDTH - 0.8), Inches(0.02))
    elif preset == 'left-half':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(bottom), Inches(SLIDE_WIDTH / 2), Inches(0.02))
    elif preset == 'right-half':
        line = master.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(SLIDE_WIDTH / 2), Inches(bottom), Inches(SLIDE_WIDTH / 2), Inches(0.02))
    else:
        return
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
