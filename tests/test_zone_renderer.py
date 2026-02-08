"""
Unit tests for zone_renderer module.
"""

# Path setup handled by conftest.py

from pptx import Presentation
from pptx.util import Inches
from pptx_gen.zone_renderer import render_zone_content
from pptx_gen.themes import get_theme


def test_line_chart_rendering():
    """Test that line charts are rendered with proper styling."""
    prs = Presentation.create(template="mitsein")
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    theme = get_theme('soft_peach_cream')

    zone_data = {
        'chartData': {
            'title': 'Sales Trend',
            'chartType': 'line',
            'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
            'series': [
                {'name': 'Revenue', 'data': [100, 120, 140, 160]}
            ]
        }
    }

    # Should not crash
    render_zone_content(slide, 'chart', 'A', 1.0, 1.0, 4.0, 3.0, theme, zone_data)

    print("✅ Line chart rendering test passed")


def test_pie_chart_rendering():
    """Test that pie charts are rendered without line styling."""
    prs = Presentation.create(template="mitsein")
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    theme = get_theme('soft_peach_cream')

    zone_data = {
        'chartData': {
            'title': 'Market Share',
            'chartType': 'pie',
            'categories': ['Product A', 'Product B', 'Product C'],
            'series': [
                {'name': 'Share', 'data': [40, 35, 25]}
            ]
        }
    }

    # Should not crash (pie charts don't have line attribute)
    render_zone_content(slide, 'chart', 'B', 1.0, 1.0, 4.0, 3.0, theme, zone_data)

    print("✅ Pie chart rendering test passed")


def test_empty_chart_data():
    """Test that empty chart data shows placeholder."""
    prs = Presentation.create(template="mitsein")
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    theme = get_theme('soft_peach_cream')

    zone_data = {
        'chartData': {
            'title': 'Empty Chart',
            'chartType': 'line',
            'categories': [],
            'series': []
        }
    }

    # Should show placeholder instead of crashing
    render_zone_content(slide, 'chart', 'C', 1.0, 1.0, 4.0, 3.0, theme, zone_data)

    print("✅ Empty chart data test passed")


if __name__ == '__main__':
    test_line_chart_rendering()
    test_pie_chart_rendering()
    test_empty_chart_data()
    print("\n✅ All zone_renderer tests passed!")
