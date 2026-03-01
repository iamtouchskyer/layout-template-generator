"""
Grid/content slide generator.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ..themes import hex_to_rgb
from ..layout import calculate_zone_positions, render_zone_content
from ..dimensions import px_to_inches_x, px_to_inches_y
from ..utils import get_blank_layout, remove_slide_placeholders

# Layout constants
PILL_WIDTH = Inches(0.8)
PILL_HEIGHT = Inches(0.3)
PILL_SPACING = Inches(0.15)
PILL_FONT_SIZE = Pt(10)
TITLE_FONT_SIZE = Pt(32)
SOURCE_FONT_SIZE = Pt(11)
SOURCE_HEIGHT = Inches(0.25)

DEFAULT_SLIDE_WIDTH_PX = 1280.0
DEFAULT_SLIDE_HEIGHT_PX = 720.0
DEFAULT_MARGINS_PX = {"top": 20.0, "right": 40.0, "bottom": 40.0, "left": 40.0}
DEFAULT_BODY_GAP_PX = 8.0
HEADER_HEIGHT_PX = {"compact": 60.0, "normal": 75.0, "spacious": 90.0}
FOOTER_HEIGHT_PX = {"compact": 15.0, "normal": 20.0, "spacious": 25.0}


def _to_float(value, default):
    try:
        return float(value)
    except (TypeError, ValueError):
        return float(default)


def _read_bounds(bounds: dict, defaults: dict) -> dict:
    raw = bounds if isinstance(bounds, dict) else {}
    return {
        "x": _to_float(raw.get("x"), defaults["x"]),
        "y": _to_float(raw.get("y"), defaults["y"]),
        "width": _to_float(raw.get("width"), defaults["width"]),
        "height": _to_float(raw.get("height"), defaults["height"]),
    }


def resolve_content_shell_geometry(config: dict) -> dict:
    """Resolve shared content-shell geometry for content-grid and smartart slides.

    Uses frontend-exported bounds when present, and computes dynamic defaults from
    current slide ratio when bounds are missing.
    """
    slide_cfg = config.get("slide", {}) if isinstance(config.get("slide"), dict) else {}
    slide_w = _to_float(slide_cfg.get("width"), DEFAULT_SLIDE_WIDTH_PX)
    slide_h = _to_float(slide_cfg.get("height"), DEFAULT_SLIDE_HEIGHT_PX)

    base_margin_cfg = slide_cfg.get("baseMargin", {}) if isinstance(slide_cfg.get("baseMargin"), dict) else {}
    margins = {
        "top": _to_float(base_margin_cfg.get("top"), DEFAULT_MARGINS_PX["top"]),
        "right": _to_float(base_margin_cfg.get("right"), DEFAULT_MARGINS_PX["right"]),
        "bottom": _to_float(base_margin_cfg.get("bottom"), DEFAULT_MARGINS_PX["bottom"]),
        "left": _to_float(base_margin_cfg.get("left"), DEFAULT_MARGINS_PX["left"]),
    }

    slide_master = config.get("slideMaster", {}) if isinstance(config.get("slideMaster"), dict) else {}
    content_areas = slide_master.get("contentAreas", {})
    if not isinstance(content_areas, dict):
        content_areas = {}

    title_style = str(content_areas.get("titleStyle", "with-tag"))
    source_style = str(content_areas.get("sourceStyle", "citation"))
    has_title = title_style != "none"
    has_source = source_style == "citation"

    header_style = str(content_areas.get("headerHeight", "compact"))
    footer_style = str(content_areas.get("footerHeight", "compact"))
    header_h = HEADER_HEIGHT_PX.get(header_style, HEADER_HEIGHT_PX["compact"])
    footer_h = FOOTER_HEIGHT_PX.get(footer_style, FOOTER_HEIGHT_PX["compact"])
    body_gap = _to_float(content_areas.get("bodyGap"), DEFAULT_BODY_GAP_PX)

    default_header = {
        "x": margins["left"],
        "y": margins["top"],
        "width": max(0.0, slide_w - margins["left"] - margins["right"]),
        "height": header_h,
    }
    body_y = margins["top"] + header_h + body_gap if has_title else margins["top"]
    default_body = {
        "x": margins["left"],
        "y": body_y,
        "width": max(0.0, slide_w - margins["left"] - margins["right"]),
        "height": max(0.0, slide_h - body_y - margins["bottom"]),
    }
    default_footer = {
        "x": margins["left"],
        "y": max(0.0, slide_h - margins["bottom"]),
        "width": max(0.0, slide_w - margins["left"] - margins["right"]),
        "height": footer_h,
    }

    header_px = _read_bounds(content_areas.get("headerBounds"), default_header)
    body_px = _read_bounds(content_areas.get("bodyBounds"), default_body)
    footer_px = _read_bounds(content_areas.get("footerBounds"), default_footer)

    content_title = str(config.get("contentTitle") or "市场趋势分析")
    content_tag = str(config.get("contentTag") or "分析报告")
    content_source = str(config.get("contentSource") or "行业研究报告 2024")

    return {
        "title_style": title_style,
        "source_style": source_style,
        "has_title": has_title,
        "has_source": has_source,
        "header_px": header_px,
        "body_px": body_px,
        "footer_px": footer_px,
        "header_in": {
            "x": px_to_inches_x(header_px["x"]),
            "y": px_to_inches_y(header_px["y"]),
            "w": px_to_inches_x(header_px["width"]),
            "h": px_to_inches_y(header_px["height"]),
        },
        "body_in": {
            "x": px_to_inches_x(body_px["x"]),
            "y": px_to_inches_y(body_px["y"]),
            "w": px_to_inches_x(body_px["width"]),
            "h": px_to_inches_y(body_px["height"]),
        },
        "footer_in": {
            "x": px_to_inches_x(footer_px["x"]),
            "y": px_to_inches_y(footer_px["y"]),
            "w": px_to_inches_x(footer_px["width"]),
            "h": px_to_inches_y(footer_px["height"]),
        },
        "content_title": content_title,
        "content_tag": content_tag,
        "content_source": content_source,
    }


def generate_grid_slide(prs, config, theme):
    """Generate grid/content slide with dynamic layout from config.

    Uses slide layout with placeholders to inherit from slide master.
    Uses pre-calculated bounds from frontend for WYSIWYG accuracy.
    """
    shell = resolve_content_shell_geometry(config)
    title_style = shell["title_style"]
    has_title = shell["has_title"]

    # Always start from blank layout to avoid accidental template placeholders
    # like "Click to add title/object" leaking into exported slides.
    slide = prs.slides.add_slide(get_blank_layout(prs))
    remove_slide_placeholders(slide)

    # Get grid config
    grid_config = config.get('grid', {})
    layout_id = grid_config.get('layout', 'two-col-equal')
    zones = grid_config.get('zones', [])

    body_x = shell["body_in"]["x"]
    body_y = shell["body_in"]["y"]
    body_w = shell["body_in"]["w"]
    body_h = shell["body_in"]["h"]
    header_x = shell["header_in"]["x"]
    header_y = shell["header_in"]["y"]
    header_w = shell["header_in"]["w"]
    header_h = shell["header_in"]["h"]
    footer_x = shell["footer_in"]["x"]
    footer_y = shell["footer_in"]["y"]
    footer_w = shell["footer_in"]["w"]

    ZONE_GAP = px_to_inches_x(24)  # Gap between zones (24px)

    content_title = shell["content_title"]
    content_tag = shell["content_tag"]
    content_source = shell["content_source"]

    # Render title directly into shell header bounds.
    if has_title:
        set_title_with_style(
            slide,
            None,
            content_title,
            content_tag,
            theme,
            title_style,
            bounds=(header_x, header_y, header_w, header_h),
        )

    # Add source citation if enabled
    has_source = shell["has_source"]
    if has_source:
        add_source_citation_dynamic(slide, content_source, theme,
                                     x=footer_x, y=footer_y, width=footer_w)

    # Calculate zone positions based on layout (using body bounds)
    zone_positions = calculate_zone_positions(layout_id, zones, body_x, body_y, body_w, body_h, ZONE_GAP)

    # Build zone data lookup for chart/content data
    zone_data_map = {z['id']: z for z in zones}

    # Render each zone
    for zone_pos in zone_positions:
        zone_id = zone_pos['id']
        content_type = zone_pos.get('content', 'text')
        x, y, w, h = zone_pos['x'], zone_pos['y'], zone_pos['w'], zone_pos['h']

        # Get full zone data (includes chartData, textData, etc.)
        zone_data = zone_data_map.get(zone_id, {})

        render_zone_content(slide, content_type, zone_id, x, y, w, h, theme, zone_data)

    return slide


def set_title_with_style(
    slide,
    title_placeholder,
    title: str,
    tag: str,
    theme: dict,
    title_style: str,
    bounds: tuple[float, float, float, float] | None = None,
):
    """Set title with optional pill-shaped tag.

    For 'with-tag' style, adds a rounded rectangle pill before the title.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    text_color = hex_to_rgb(theme['primary'])
    accent_color = hex_to_rgb(theme['accent'])

    if bounds is not None:
        bx, by, bw, bh = bounds
        ph_left = Inches(bx)
        ph_top = Inches(by)
        ph_width = Inches(bw)
        ph_height = Inches(bh)
    elif title_placeholder is not None:
        ph_left = title_placeholder.left
        ph_top = title_placeholder.top
        ph_width = title_placeholder.width
        ph_height = title_placeholder.height
    else:
        # Fallback to default header area
        ph_left = Inches(0.42)
        ph_top = Inches(0.2)
        ph_width = Inches(12.5)
        ph_height = Inches(0.62)

    tf = title_placeholder.text_frame if title_placeholder is not None else None
    if tf is not None:
        tf.clear()

    if title_style == 'with-tag' and tag:
        # Add pill-shaped tag as a separate shape
        pill_left = ph_left
        pill_top = ph_top + (ph_height - PILL_HEIGHT) // 2

        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            pill_left, pill_top, PILL_WIDTH, PILL_HEIGHT
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = accent_color
        pill.line.fill.background()

        # Add tag text in pill
        pill.text_frame.paragraphs[0].text = tag
        pill.text_frame.paragraphs[0].font.size = PILL_FONT_SIZE
        pill.text_frame.paragraphs[0].font.bold = True
        pill.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Offset title to the right of the pill
        title_left = pill_left + PILL_WIDTH + PILL_SPACING
        title_width = ph_width - PILL_WIDTH - PILL_SPACING

        # Create title textbox (don't use placeholder since we need custom positioning)
        title_box = slide.shapes.add_textbox(title_left, ph_top, title_width, ph_height)
        title_box.text_frame.word_wrap = True
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = TITLE_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = text_color
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Clear the original placeholder (we're using custom shapes)
        if tf is not None:
            tf.paragraphs[0].text = ""
    else:
        # Use custom text box so geometry is always controlled by headerBounds.
        title_box = slide.shapes.add_textbox(ph_left, ph_top, ph_width, ph_height)
        title_box.text_frame.word_wrap = True
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = TITLE_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = text_color
        if tf is not None:
            tf.paragraphs[0].text = ""


def add_source_citation_dynamic(slide, source: str, theme: dict, x: float, y: float, width: float):
    """Add source citation at specified position."""
    text_box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), SOURCE_HEIGHT
    )
    p = text_box.text_frame.paragraphs[0]
    p.text = f"数据来源：{source}"
    p.font.size = SOURCE_FONT_SIZE
    p.font.color.rgb = hex_to_rgb(theme.get('text_muted', '#888888'))
