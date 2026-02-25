"""
Cover slide generator with multiple layout support.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..themes import hex_to_rgb
from ..dimensions import px_to_inches_x, px_to_inches_y
from ..cover_layouts import get_cover_layout, COVER_LAYOUTS
from ..utils import get_blank_layout, remove_slide_placeholders


def generate_cover_slide(prs, config, theme):
    """Generate cover slide with specified layout.

    Args:
        prs: Presentation object
        config: Config dict, may contain 'coverLayout' key
        theme: Theme colors dict
    """
    slide = prs.slides.add_slide(get_blank_layout(prs))
    remove_slide_placeholders(slide)

    # Get layout from config or use default
    layout_id = config.get('coverLayout', 'cross_rectangles')
    layout = get_cover_layout(layout_id)

    # Get content from config
    cover_content = config.get('coverContent', {})
    year = cover_content.get('year', '2025')
    title = cover_content.get('title', '年度汇报')
    highlight = cover_content.get('highlight', '')
    subtitle = cover_content.get('subtitle', '在这里输入演示文稿的副标题或简要描述')
    tag = cover_content.get('tag', '')
    brand_tag = cover_content.get('brandTag', '')

    # Footer content
    footer_content = cover_content.get('footer', {})
    location = footer_content.get('location', '芝士科技大厦')
    date = footer_content.get('date', '2025.01')
    contact = footer_content.get('contact', '400-123-4567')
    logo_text = footer_content.get('logo', 'LOGO')

    # Render shapes
    for shape_config in layout.get('shapes', []):
        _render_shape(slide, shape_config, theme)

    # Render text areas
    text_areas = layout.get('text_areas', {})

    if 'year' in text_areas and year:
        _render_text_area(slide, text_areas['year'], year, theme)

    if 'title' in text_areas and title:
        _render_text_area(slide, text_areas['title'], title, theme)

    if 'highlight' in text_areas and highlight:
        _render_text_area(slide, text_areas['highlight'], highlight, theme)

    if 'subtitle' in text_areas and subtitle:
        _render_text_area(slide, text_areas['subtitle'], subtitle, theme)

    if 'tag' in text_areas and tag:
        _render_text_area(slide, text_areas['tag'], tag, theme)

    # Render brand tag if present
    if 'brand_tag' in layout and brand_tag:
        _render_text_area(slide, layout['brand_tag'], brand_tag, theme)

    # Render footer
    footer_config = layout.get('footer', {})
    if footer_config.get('enabled', False):
        _render_footer(slide, footer_config, theme, location, date, contact, logo_text)

    return slide


def _resolve_color(color_ref: str, theme: dict) -> RGBColor:
    """Resolve color reference to RGBColor.

    Supports:
    - Hex colors: '#FFFFFF'
    - Theme references: 'accent', 'accent1', 'accent2', 'text', 'bg', 'bg2', 'tx2', 'text_muted', 'primary'
    """
    if color_ref.startswith('#'):
        return hex_to_rgb(color_ref)

    # Map color references to theme keys
    color_map = {
        'accent': 'accent',
        'accent1': 'accent1',
        'accent2': 'accent2',
        'accent3': 'accent3',
        'accent4': 'accent4',
        'accent5': 'accent5',
        'accent6': 'accent6',
        'text': 'text',
        'text_muted': 'text_muted',
        'primary': 'primary',
        'bg': 'bg',
        'bg2': 'card_bg',  # bg2 maps to card_bg (light background)
        'tx2': 'accent1',  # tx2 maps to accent1 (dark text/shape color)
    }

    theme_key = color_map.get(color_ref, 'accent')

    # Try to get from theme, with fallbacks
    if theme_key in theme:
        return hex_to_rgb(theme[theme_key])

    # Fallback for accent1-6 if not defined
    if theme_key.startswith('accent') and theme_key != 'accent':
        return hex_to_rgb(theme.get('accent', '#1C58A0'))

    # Fallback for text_muted
    if theme_key == 'text_muted':
        return hex_to_rgb(theme.get('text_muted', '#888888'))

    return hex_to_rgb(theme.get('accent', '#1C58A0'))


def _render_shape(slide, shape_config: dict, theme: dict):
    """Render a decorative shape."""
    shape_type = shape_config.get('type', 'rect')
    x = px_to_inches_x(shape_config.get('x', 0))
    y = px_to_inches_y(shape_config.get('y', 0))
    w = px_to_inches_x(shape_config.get('w', 100))
    h = px_to_inches_y(shape_config.get('h', 100))
    fill = shape_config.get('fill', 'accent')
    rotation = shape_config.get('rotation', 0)

    # Map shape type to MSO_SHAPE
    shape_type_map = {
        'rect': MSO_SHAPE.RECTANGLE,
        'roundRect': MSO_SHAPE.ROUNDED_RECTANGLE,
        'ellipse': MSO_SHAPE.OVAL,
        'rtTriangle': MSO_SHAPE.RIGHT_TRIANGLE,
    }
    mso_shape = shape_type_map.get(shape_type, MSO_SHAPE.RECTANGLE)

    shape = slide.shapes.add_shape(mso_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _resolve_color(fill, theme)
    shape.line.fill.background()

    if rotation:
        shape.rotation = rotation


def _render_text_area(slide, text_config: dict, text: str, theme: dict):
    """Render a text area."""
    x = px_to_inches_x(text_config.get('x', 0))
    y = px_to_inches_y(text_config.get('y', 0))
    w = px_to_inches_x(text_config.get('w', 400))
    h = px_to_inches_y(text_config.get('h', 50))
    font_size = text_config.get('font_size', 24)
    color = text_config.get('color', 'text')
    bold = text_config.get('bold', False)
    align = text_config.get('align', 'left')

    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_box.text_frame.word_wrap = True

    p = text_box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _resolve_color(color, theme)

    # Set alignment
    align_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
    }
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)


def _render_footer(slide, footer_config: dict, theme: dict,
                   location: str, date: str, contact: str, logo_text: str):
    """Render footer bar with icons and text."""
    bar_config = footer_config.get('bar', {})

    # Render footer bar background
    x = px_to_inches_x(bar_config.get('x', 38))
    y = px_to_inches_y(bar_config.get('y', 613))
    w = px_to_inches_x(bar_config.get('w', 1203))
    h = px_to_inches_y(bar_config.get('h', 65))
    bar_type = bar_config.get('type', 'roundRect')
    bar_fill = bar_config.get('fill', 'bg2')

    shape_type_map = {
        'rect': MSO_SHAPE.RECTANGLE,
        'roundRect': MSO_SHAPE.ROUNDED_RECTANGLE,
    }
    mso_shape = shape_type_map.get(bar_type, MSO_SHAPE.ROUNDED_RECTANGLE)

    bar = slide.shapes.add_shape(mso_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _resolve_color(bar_fill, theme)
    bar.line.fill.background()

    # Render footer items
    items = footer_config.get('items', [])
    content_map = {
        '芝士科技大厦': location,
        '2025.01': date,
        '400-123-4567': contact,
        'LOGO': logo_text,
    }

    for item in items:
        item_type = item.get('type', 'text')

        if item_type == 'text':
            ix = px_to_inches_x(item.get('x', 0))
            iy = px_to_inches_y(item.get('y', 0))
            default_text = item.get('text', '')
            actual_text = content_map.get(default_text, default_text)
            color = item.get('color', 'text')

            text_box = slide.shapes.add_textbox(Inches(ix), Inches(iy), Inches(2), Inches(0.3))
            p = text_box.text_frame.paragraphs[0]
            p.text = actual_text
            p.font.size = Pt(11)
            p.font.color.rgb = _resolve_color(color, theme)

        elif item_type == 'icon':
            # Render icon as small shape placeholder
            ix = px_to_inches_x(item.get('x', 0))
            iy = px_to_inches_y(item.get('y', 0))
            color = item.get('color', 'accent')
            icon_name = item.get('icon', 'location')

            # Use small oval as icon placeholder
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(ix), Inches(iy), Inches(0.2), Inches(0.2)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = _resolve_color(color, theme)
            icon.line.fill.background()
