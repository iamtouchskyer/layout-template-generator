"""
Master placeholders (logo, date, page-number) for slide master.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from .themes import hex_to_rgb
from .dimensions import px_to_inches_x, px_to_inches_y, get_pptx_dimensions


def add_master_placeholders(master, placeholders: dict, theme: dict):
    """Add placeholder shapes to slide master.

    Supports placeholder format:
    {
        "logo": { "enabled": true, "position": "tl", "positionConfig": {...}, "size": "medium", "sizeConfig": {...}, "imageUrl": "..." },
        "page-number": { "enabled": true, "position": "br", "positionConfig": {...} },
        "date": { "enabled": true, "position": "bl", "positionConfig": {...} }
    }
    """
    muted_color = hex_to_rgb(theme.get('text_muted', '#888888'))
    accent_color = hex_to_rgb(theme['accent'])

    SLIDE_WIDTH, SLIDE_HEIGHT = get_pptx_dimensions()

    # Handle both old array format and new object format
    if isinstance(placeholders, list):
        placeholders = {ph_id: {'enabled': True} for ph_id in placeholders}

    enable_sldnum = False
    enable_dt = False
    sldnum_config = None
    dt_config = None

    for ph_id, ph_config in placeholders.items():
        if not ph_config.get('enabled', False):
            continue

        pos_config = ph_config.get('positionConfig', {})

        if ph_id == 'page-number':
            enable_sldnum = True
            sldnum_config = ph_config
        elif ph_id == 'date':
            enable_dt = True
            dt_config = ph_config
        elif ph_id == 'logo':
            _add_logo_placeholder(master, ph_config, pos_config, accent_color, muted_color, SLIDE_WIDTH, SLIDE_HEIGHT)

    # Enable/disable footer placeholders using new python-pptx API
    set_header_footer(master, sldNum=enable_sldnum, dt=enable_dt)

    # Set positions for enabled placeholders
    if enable_sldnum and sldnum_config:
        _set_footer_placeholder_position(master.slide_number_placeholder, sldnum_config, SLIDE_WIDTH, SLIDE_HEIGHT)

    if enable_dt and dt_config:
        _set_footer_placeholder_position(master.date_placeholder, dt_config, SLIDE_WIDTH, SLIDE_HEIGHT)


def set_header_footer(master, sldNum: bool = False, dt: bool = False, ftr: bool = False):
    """Set header/footer visibility flags on slide master.

    Uses the new python-pptx API for footer placeholder control.
    """
    master.show_slide_number(sldNum)
    master.show_date(dt)


def _set_footer_placeholder_position(placeholder, config: dict, SLIDE_WIDTH: float, SLIDE_HEIGHT: float):
    """Set position for a footer placeholder (date, slide_number, footer)."""
    if placeholder is None:
        return

    pos_config = config.get('positionConfig', {})
    position = config.get('position', 'br')

    width = px_to_inches_x(pos_config.get('width', 100))
    height = px_to_inches_y(pos_config.get('height', 30))

    if position == 'tl':
        x = px_to_inches_x(pos_config.get('left', 40))
        y = px_to_inches_y(pos_config.get('top', 20))
    elif position == 'tr':
        x = SLIDE_WIDTH - width - px_to_inches_x(pos_config.get('right', 40))
        y = px_to_inches_y(pos_config.get('top', 20))
    elif position == 'bl':
        x = px_to_inches_x(pos_config.get('left', 40))
        y = SLIDE_HEIGHT - height - px_to_inches_y(pos_config.get('bottom', 20))
    elif position == 'br':
        x = SLIDE_WIDTH - width - px_to_inches_x(pos_config.get('right', 40))
        y = SLIDE_HEIGHT - height - px_to_inches_y(pos_config.get('bottom', 20))
    else:
        x = px_to_inches_x(pos_config.get('left', 40))
        y = px_to_inches_y(pos_config.get('top', 700))

    placeholder.left = Inches(x)
    placeholder.top = Inches(y)
    placeholder.width = Inches(width)
    placeholder.height = Inches(height)


def _add_logo_placeholder(master, ph_config, pos_config, accent_color, muted_color, SLIDE_WIDTH, SLIDE_HEIGHT):
    """Add logo placeholder to slide master."""
    import base64
    import io

    size_config = ph_config.get('sizeConfig', {})
    width = px_to_inches_x(size_config.get('width', 80))
    height = px_to_inches_y(size_config.get('height', 30))

    position = ph_config.get('position', 'tl')
    if position == 'tl':
        x = px_to_inches_x(pos_config.get('left', 40))
        y = px_to_inches_y(pos_config.get('top', 24))
    elif position == 'tr':
        x = SLIDE_WIDTH - width - px_to_inches_x(pos_config.get('right', 40))
        y = px_to_inches_y(pos_config.get('top', 24))
    elif position == 'bl':
        x = px_to_inches_x(pos_config.get('left', 40))
        y = SLIDE_HEIGHT - height - px_to_inches_y(pos_config.get('bottom', 40))
    elif position == 'br':
        x = SLIDE_WIDTH - width - px_to_inches_x(pos_config.get('right', 40))
        y = SLIDE_HEIGHT - height - px_to_inches_y(pos_config.get('bottom', 40))
    else:
        x, y = 0.4, 0.25

    image_url = ph_config.get('imageUrl')
    if image_url and isinstance(image_url, str) and image_url.startswith('data:image'):
        # Validate data URI format: data:image/png;base64,<data>
        MAX_IMAGE_SIZE = 5 * 1024 * 1024  # 5MB limit
        try:
            if ',' not in image_url:
                raise ValueError("Invalid data URI format: missing comma separator")
            header, encoded = image_url.split(',', 1)
            if not header.startswith('data:image/'):
                raise ValueError(f"Invalid image type: {header}")
            if len(encoded) > MAX_IMAGE_SIZE:
                raise ValueError(f"Image data exceeds {MAX_IMAGE_SIZE} bytes limit")
            image_data = base64.b64decode(encoded, validate=True)
            image_stream = io.BytesIO(image_data)
            master.shapes.add_picture(image_stream, Inches(x), Inches(y), Inches(width), Inches(height))
        except (ValueError, base64.binascii.Error) as e:
            import logging
            logging.warning(f"Logo image decode failed: {e}")
            _add_logo_text_placeholder(master, x, y, width, height, muted_color)
    else:
        _add_logo_shape_placeholder(master, x, y, width, height, accent_color)


def _add_logo_text_placeholder(master, x, y, width, height, muted_color):
    """Add text-based logo placeholder."""
    ph = master.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    p = ph.text_frame.paragraphs[0]
    p.text = "LOGO"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = muted_color


def _add_logo_shape_placeholder(master, x, y, width, height, accent_color):
    """Add shape-based logo placeholder."""
    logo_shape = master.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    logo_shape.fill.solid()
    logo_shape.fill.fore_color.rgb = accent_color
    logo_shape.line.fill.background()
    logo_shape.text_frame.paragraphs[0].text = "LOGO"
    logo_shape.text_frame.paragraphs[0].font.size = Pt(10)
    logo_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    logo_shape.text_frame.paragraphs[0].font.bold = True
    logo_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
