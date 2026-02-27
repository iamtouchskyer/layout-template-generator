"""
Zone content rendering for PPTX generation.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .themes import hex_to_rgb


MIN_DIM_IN = 0.02


def _safe_dim(value: float, minimum: float = MIN_DIM_IN) -> float:
    """Ensure dimension value is positive for OOXML shape extents."""
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        numeric = minimum
    return numeric if numeric > minimum else minimum


def _add_shape(slide, shape_type, x, y, w, h):
    return slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(_safe_dim(w)),
        Inches(_safe_dim(h)),
    )


def _add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(_safe_dim(w)),
        Inches(_safe_dim(h)),
    )


def _add_chart(slide, chart_type, x, y, w, h, chart_data):
    return slide.shapes.add_chart(
        chart_type,
        Inches(x),
        Inches(y),
        Inches(_safe_dim(w, 0.1)),
        Inches(_safe_dim(h, 0.1)),
        chart_data,
    )


def render_zone_content(slide, content_type: str, zone_id: str, x: float, y: float,
                        w: float, h: float, theme: dict, zone_data: dict = None):
    """Render content in a zone based on content type.

    Args:
        slide: PowerPoint slide object
        content_type: Type of content (chart, image, metric, table, bullets, text)
        zone_id: Zone identifier (A, B, C, D, etc.)
        x, y: Zone position in inches
        w, h: Zone dimensions in inches
        theme: Theme colors dict
        zone_data: Full zone data including chartData, textData, etc.
    """
    zone_data = zone_data or {}
    accent_color = hex_to_rgb(theme.get('accent', '#1C58A0'))
    text_color = hex_to_rgb(theme.get('text', '#333333'))
    card_bg = hex_to_rgb('#FFFFFF')
    card_border = hex_to_rgb(theme.get('card_border', '#E8E8E8'))
    muted_color = hex_to_rgb(theme.get('text_muted', '#888888'))

    # Add card background
    card = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = card_bg
    card.line.color.rgb = card_border

    # Content padding
    padding = min(0.15, max(0.02, min(_safe_dim(w), _safe_dim(h)) * 0.18))
    content_x = x + padding
    content_y = y + padding
    content_w = _safe_dim(w - 2 * padding, 0.05)
    content_h = _safe_dim(h - 2 * padding, 0.05)

    if _is_compact_cell(content_w, content_h):
        _render_compact_zone(
            slide,
            zone_id,
            content_type,
            content_x,
            content_y,
            content_w,
            content_h,
            text_color,
            muted_color,
            zone_data,
        )
        return

    if content_type == 'chart':
        chart_data = zone_data.get('chartData', {})
        _render_chart_zone(slide, zone_id, content_x, content_y, content_w, content_h, text_color, theme, chart_data)
    elif content_type == 'image':
        _render_image_zone(slide, zone_id, content_x, content_y, content_w, content_h, x, y, w, h, text_color, card_border)
    elif content_type == 'metric':
        _render_metric_zone(slide, zone_id, content_x, content_y, content_w, content_h, accent_color, muted_color)
    elif content_type == 'table':
        _render_table_zone(slide, zone_id, content_x, content_y, content_w, content_h, x, y, w, h, text_color, card_border)
    elif content_type == 'bullets':
        _render_bullets_zone(slide, zone_id, content_x, content_y, content_w, content_h, y, h, text_color, accent_color)
    else:
        _render_text_zone(slide, zone_id, content_x, content_y, content_w, content_h, text_color, zone_data)


def _render_chart_zone(slide, zone_id, content_x, content_y, content_w, content_h, text_color, theme, chart_data):
    """Render real chart using python-pptx chart API."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    # Chart title
    title = chart_data.get('title', f'图表 {zone_id}')
    header = _add_textbox(slide, content_x, content_y, content_w, 0.35)
    p = header.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color

    # Chart area position (below title)
    chart_x = content_x
    chart_y = content_y + 0.4
    chart_w = content_w
    chart_h = content_h - 0.5

    # If no chart data, show placeholder
    if not chart_data.get('categories') or not chart_data.get('series'):
        _render_chart_placeholder(slide, chart_x, chart_y, chart_w, chart_h)
        return

    # Build chart data
    categories = chart_data.get('categories', [])
    series_list = chart_data.get('series', [])
    chart_type_str = chart_data.get('chartType', 'line')

    # Map chart type string to XL_CHART_TYPE
    chart_type_map = {
        'line': XL_CHART_TYPE.LINE,
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'pie': XL_CHART_TYPE.PIE,
        'area': XL_CHART_TYPE.AREA,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
    }
    xl_chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.LINE)

    # Create chart data
    cd = CategoryChartData()
    cd.categories = categories

    for series in series_list:
        series_name = series.get('name', 'Series')
        series_values = series.get('data', [])
        cd.add_series(series_name, series_values)

    # Add chart to slide
    chart_shape = _add_chart(slide, xl_chart_type, chart_x, chart_y, chart_w, chart_h, cd)

    # Style the chart
    chart = chart_shape.chart
    chart.has_legend = False

    # Set line color to accent color for line charts
    if chart_type_str == 'line' and chart.series:
        accent_color = hex_to_rgb(theme['accent'])
        for series in chart.series:
            if hasattr(series.format, 'line'):
                series.format.line.color.rgb = accent_color
                series.format.line.width = Pt(2)


def _render_chart_placeholder(slide, x, y, w, h):
    """Render chart placeholder when no data available."""
    area = _add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h)
    area.fill.solid()
    area.fill.fore_color.rgb = hex_to_rgb('#F5F5F5')
    area.line.color.rgb = hex_to_rgb('#E0E0E0')

    icon = _add_textbox(slide, x + w/2 - 0.3, y + h/2 - 0.2, 0.6, 0.4)
    icon.text_frame.paragraphs[0].text = "📊"
    icon.text_frame.paragraphs[0].font.size = Pt(24)
    icon.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _render_image_zone(slide, zone_id, content_x, content_y, content_w, content_h, x, y, w, h, text_color, card_border):
    """Render image placeholder zone."""
    header = _add_textbox(slide, content_x, content_y, content_w, 0.35)
    p = header.text_frame.paragraphs[0]
    p.text = f"图片 {zone_id}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color

    img_area = _add_shape(slide, MSO_SHAPE.RECTANGLE, content_x + 0.2, content_y + 0.5, content_w - 0.4, content_h - 0.7)
    img_area.fill.solid()
    img_area.fill.fore_color.rgb = hex_to_rgb('#E8E8E8')
    img_area.line.color.rgb = card_border

    icon = _add_textbox(slide, x + w/2 - 0.3, y + h/2 - 0.2, 0.6, 0.4)
    icon.text_frame.paragraphs[0].text = "🖼️"
    icon.text_frame.paragraphs[0].font.size = Pt(24)
    icon.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _render_metric_zone(slide, zone_id, content_x, content_y, content_w, content_h, accent_color, muted_color):
    """Render metric display zone."""
    metric_box = _add_textbox(slide, content_x, content_y + content_h * 0.2, content_w, content_h * 0.4)
    metric_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = metric_box.text_frame.paragraphs[0]
    p.text = "85%"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = accent_color

    label_box = _add_textbox(slide, content_x, content_y + content_h * 0.6, content_w, 0.4)
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = label_box.text_frame.paragraphs[0]
    p.text = f"关键指标 {zone_id}"
    p.font.size = Pt(14)
    p.font.color.rgb = muted_color


def _render_table_zone(slide, zone_id, content_x, content_y, content_w, content_h, x, y, w, h, text_color, card_border):
    """Render table placeholder zone."""
    header = _add_textbox(slide, content_x, content_y, content_w, 0.35)
    p = header.text_frame.paragraphs[0]
    p.text = f"表格 {zone_id}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color

    table_area = _add_shape(slide, MSO_SHAPE.RECTANGLE, content_x + 0.1, content_y + 0.5, content_w - 0.2, content_h - 0.7)
    table_area.fill.solid()
    table_area.fill.fore_color.rgb = hex_to_rgb('#FAFAFA')
    table_area.line.color.rgb = card_border

    icon = _add_textbox(slide, x + w/2 - 0.3, y + h/2 - 0.2, 0.6, 0.4)
    icon.text_frame.paragraphs[0].text = "📋"
    icon.text_frame.paragraphs[0].font.size = Pt(24)
    icon.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _render_bullets_zone(slide, zone_id, content_x, content_y, content_w, content_h, y, h, text_color, accent_color):
    """Render bullet list zone."""
    header = _add_textbox(slide, content_x, content_y, content_w, 0.35)
    p = header.text_frame.paragraphs[0]
    p.text = f"要点 {zone_id}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color

    bullets = ["要点一：关键信息", "要点二：重要说明", "要点三：总结概括"]
    for i, bullet_text in enumerate(bullets):
        if content_y + 0.55 + i * 0.45 > y + h - 0.3:
            break

        dot = _add_shape(slide, MSO_SHAPE.OVAL, content_x + 0.05, content_y + 0.55 + i * 0.45, 0.08, 0.08)
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent_color
        dot.line.fill.background()

        text = _add_textbox(slide, content_x + 0.2, content_y + 0.45 + i * 0.45, content_w - 0.3, 0.4)
        text.text_frame.word_wrap = True
        p = text.text_frame.paragraphs[0]
        p.text = bullet_text
        p.font.size = Pt(11)
        p.font.color.rgb = text_color


def _render_text_zone(slide, zone_id, content_x, content_y, content_w, content_h, text_color, zone_data=None):
    """Render text content zone."""
    zone_data = zone_data or {}
    text_data = zone_data.get('textData', {})

    # Use provided title or default
    title = text_data.get('title', f"区域 {zone_id}")
    body_text = text_data.get('body', "这里是文本内容区域，可以放置段落、描述或说明文字。")

    header = _add_textbox(slide, content_x, content_y, content_w, 0.35)
    p = header.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = text_color

    body = _add_textbox(slide, content_x, content_y + 0.45, content_w, content_h - 0.6)
    body.text_frame.word_wrap = True
    p = body.text_frame.paragraphs[0]
    p.text = body_text
    p.font.size = Pt(11)
    p.font.color.rgb = text_color


def _is_compact_cell(content_w: float, content_h: float) -> bool:
    """Whether zone cell is too small for full card rendering."""
    return content_w < 1.25 or content_h < 1.0


def _render_compact_zone(slide, zone_id, content_type, x, y, w, h, text_color, muted_color, zone_data=None):
    """Render compact cell while preserving real UI text payload."""
    zone_data = zone_data or {}
    raw_title_text, raw_body_text = _compact_text_payload(zone_id, content_type, zone_data)
    title_size, subtitle_size = _compact_font_sizes(w, h)
    title_lines, body_lines = _compact_line_budget(w, h)
    title_chars = _compact_chars_per_line(w, title_size)
    body_chars = _compact_chars_per_line(w, subtitle_size)

    title_text = _truncate_compact_text(raw_title_text, title_lines, title_chars)
    body_text = _truncate_compact_text(raw_body_text, body_lines, body_chars)

    title_h = _compact_title_height(h, title_size, title_lines, has_body=bool(body_text))
    title = _add_textbox(slide, x, y, w, title_h)
    title.text_frame.word_wrap = True
    title.text_frame.margin_left = 0
    title.text_frame.margin_right = 0
    title.text_frame.margin_top = 0
    title.text_frame.margin_bottom = 0
    title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    if not body_text:
        return

    body_gap = min(0.04, h * 0.08)
    body_y = y + title_h + body_gap
    body_h = max(0.1, h - title_h - body_gap)

    subtitle = _add_textbox(slide, x, body_y, w, body_h)
    subtitle.text_frame.word_wrap = True
    subtitle.text_frame.margin_left = 0
    subtitle.text_frame.margin_right = 0
    subtitle.text_frame.margin_top = 0
    subtitle.text_frame.margin_bottom = 0
    subtitle.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = subtitle.text_frame.paragraphs[0]
    p.text = body_text
    p.font.size = Pt(subtitle_size)
    p.font.color.rgb = muted_color
    p.alignment = PP_ALIGN.CENTER


def _compact_font_sizes(w: float, h: float) -> tuple[int, int]:
    """Adaptive font sizes for dense compact cells."""
    cell = min(_safe_dim(w), _safe_dim(h))
    if cell < 0.35:
        return 5, 4
    if cell < 0.5:
        return 6, 5
    if cell < 0.8:
        return 7, 6
    return 8, 7


def _compact_line_budget(w: float, h: float) -> tuple[int, int]:
    """Estimate max lines for (title, body) in compact mode."""
    cell = min(_safe_dim(w), _safe_dim(h))
    if h < 0.4 or cell < 0.28:
        return 1, 0
    if h < 0.62 or w < 0.62:
        return 1, 1
    if h < 0.9:
        return 2, 1
    return 2, 2


def _compact_chars_per_line(w: float, font_size: int) -> int:
    """Approximate character capacity per line for compact text box width."""
    width_pts = _safe_dim(w) * 72.0
    avg_char_width = max(4.0, font_size * 0.9)
    chars = int(width_pts / avg_char_width)
    return max(2, min(28, chars))


def _truncate_compact_text(text: str, max_lines: int, chars_per_line: int) -> str:
    """Clamp compact text length and append ellipsis when overflow."""
    if max_lines <= 0:
        return ""
    value = str(text or "").strip()
    if not value:
        return ""
    cap = max_lines * max(1, chars_per_line)
    if len(value) <= cap:
        return value
    if cap <= 3:
        return value[:cap]
    return f"{value[: cap - 3].rstrip()}..."


def _compact_title_height(h: float, font_size: int, max_lines: int, has_body: bool) -> float:
    """Reserve enough vertical room for title, keeping body visible when present."""
    line_height = max(0.1, (font_size * 1.3) / 72.0)
    text_h = max(0.12, line_height * max(1, max_lines))
    if not has_body:
        return max(0.12, min(h, text_h + 0.04))
    return max(0.12, min(h * 0.7, max(h * 0.45, text_h + 0.02)))


def _compact_text_payload(zone_id: str, content_type: str, zone_data: dict) -> tuple[str, str]:
    """Extract meaningful compact text from zone payload."""
    if content_type == 'text':
        text_data = zone_data.get('textData', {})
        title = str(text_data.get('title') or zone_id)
        body = str(text_data.get('body') or '')
        return title, body

    if content_type == 'chart':
        chart_data = zone_data.get('chartData', {})
        title = str(chart_data.get('title') or f'图表 {zone_id}')
        series = chart_data.get('series') or []
        if isinstance(series, list) and len(series) > 0:
            first = series[0] if isinstance(series[0], dict) else {}
            series_name = str(first.get('name') or '').strip()
            if series_name:
                return title, series_name
        return title, 'Chart'

    if content_type == 'bullets':
        return f'要点 {zone_id}', '• 关键要点'
    if content_type == 'table':
        return f'表格 {zone_id}', '数据表'
    if content_type == 'image':
        return f'图片 {zone_id}', '图像内容'
    if content_type == 'metric':
        return f'指标 {zone_id}', '关键指标'

    return str(zone_id), str(content_type)
