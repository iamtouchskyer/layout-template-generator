"""Regression tests for slide layout selection in PPTX export."""

from __future__ import annotations

import os
import tempfile

from pptx import Presentation
from pptx.util import Inches

from pptx_gen import generate_pptx


CONTENT_PLACEHOLDER_TYPES = {
    "TITLE (1)",
    "CENTER_TITLE (3)",
    "SUBTITLE (4)",
    "BODY (2)",
    "OBJECT (7)",
    "CHART (8)",
    "PICTURE (18)",
}

FOOTER_PLACEHOLDER_TYPES = {
    "DATE (16)",
    "FOOTER (15)",
    "SLIDE_NUMBER (13)",
}


def _placeholder_types(slide) -> set[str]:
    result = set()
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            result.add(str(shape.placeholder_format.type))
        except Exception:
            result.add("UNKNOWN")
    return result


def _find_shape_by_exact_text(slide, text: str):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text_frame is None:
            continue
        if shape.text_frame.text.strip() == text:
            return shape
    return None


def test_cover_and_smartart_slides_do_not_include_content_placeholders():
    config = {
        "schemaVersion": 2,
        "master": {
            "theme": "forest_green",
            "masterShapes": [],
            "masterPlaceholders": {},
            "masterContentAreas": {},
        },
        "pages": [
            {
                "id": "p-cover",
                "type": "cover",
                "data": {
                    "coverLayout": "cross_rectangles",
                    "coverContent": {"title": "年度汇报"},
                },
            },
            {
                "id": "p-smartart",
                "type": "content-smartart",
                "data": {
                    "smartartType": "pyramid",
                    "smartartPlacement": "left-desc",
                    "smartartColorScheme": "colorful1",
                    "smartartItemsByType": {
                        "pyramid": [
                            {"text": "愿景", "children": [{"text": "长期目标"}]},
                            {"text": "战略", "children": [{"text": "中期规划"}]},
                            {"text": "战术", "children": [{"text": "短期行动"}]},
                            {"text": "执行", "children": [{"text": "日常任务"}]},
                        ]
                    },
                },
            },
        ],
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fp:
        output_path = fp.name

    try:
        generate_pptx(config, output_path)
        prs = Presentation.open(output_path)
        cover_ph = _placeholder_types(prs.slides[0])
        smartart_ph = _placeholder_types(prs.slides[1])
        assert cover_ph.isdisjoint(CONTENT_PLACEHOLDER_TYPES)
        assert smartart_ph.isdisjoint(CONTENT_PLACEHOLDER_TYPES)
        assert cover_ph.isdisjoint(FOOTER_PLACEHOLDER_TYPES)
        assert smartart_ph.isdisjoint(FOOTER_PLACEHOLDER_TYPES)
    finally:
        os.unlink(output_path)


def test_grid_slide_does_not_leave_template_placeholders():
    config = {
        "schemaVersion": 2,
        "master": {
            "theme": "forest_green",
            "masterShapes": [],
            "masterPlaceholders": {},
            "masterContentAreas": {
                "titleStyle": "with-tag",
                "sourceStyle": "citation",
            },
        },
        "pages": [
            {
                "id": "p-grid",
                "type": "content-grid",
                "data": {
                    "grid": {
                        "layout": "two-col-equal",
                        "zones": [
                            {"id": "A", "content": "text", "flex": 1},
                            {"id": "B", "content": "text", "flex": 1},
                        ],
                    }
                },
            }
        ],
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fp:
        output_path = fp.name

    try:
        generate_pptx(config, output_path)
        prs = Presentation.open(output_path)
        ph_types = _placeholder_types(prs.slides[0])
        assert ph_types.isdisjoint(CONTENT_PLACEHOLDER_TYPES)
        assert "DATE (16)" not in ph_types
        assert "FOOTER (15)" not in ph_types
        assert "SLIDE_NUMBER (13)" not in ph_types
    finally:
        os.unlink(output_path)


def test_divider_cards_layout_uses_full_slide_width_not_4_3_box():
    config = {
        "schemaVersion": 2,
        "master": {
            "theme": "forest_green",
            "masterShapes": [],
            "masterPlaceholders": {},
            "masterContentAreas": {},
        },
        "pages": [
            {
                "id": "p-divider",
                "type": "divider",
                "data": {
                    "divider": {
                        "layout": "cards-highlight",
                        "sectionCount": 4,
                        "sectionIndex": 0,
                        "numberStyle": "arabic",
                        "textLevel": "full",
                        "bgStyle": "solid",
                    }
                },
            }
        ],
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fp:
        output_path = fp.name

    try:
        generate_pptx(config, output_path)
        prs = Presentation.open(output_path)
        slide = prs.slides[0]
        slide_width = prs.slide_width
        max_right = max(shape.left + shape.width for shape in slide.shapes)
        # Ensure divider content stretches close to right edge on 16:9 slide.
        assert max_right > int(slide_width * 0.9)
    finally:
        os.unlink(output_path)


def test_grid_title_uses_header_bounds_position():
    config = {
        "schemaVersion": 2,
        "master": {
            "theme": "forest_green",
            "masterShapes": [],
            "masterPlaceholders": {},
            "masterContentAreas": {
                "titleStyle": "with-tag",
                "sourceStyle": "citation",
                "headerBounds": {"x": 40, "y": 20, "width": 1200, "height": 60},
            },
        },
        "pages": [
            {
                "id": "p-grid",
                "type": "content-grid",
                "data": {
                    "grid": {
                        "layout": "single",
                        "zones": [{"id": "A", "content": "text", "flex": 1}],
                    }
                },
            }
        ],
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fp:
        output_path = fp.name

    try:
        generate_pptx(config, output_path)
        prs = Presentation.open(output_path)
        slide = prs.slides[0]
        title_shape = _find_shape_by_exact_text(slide, "市场趋势分析")
        assert title_shape is not None

        expected_top = Inches(20 * (7.5 / 720))
        # allow minor rounding differences in EMU
        assert abs(title_shape.top - expected_top) < 20000
    finally:
        os.unlink(output_path)
