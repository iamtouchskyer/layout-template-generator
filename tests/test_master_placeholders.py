"""
Unit tests for master_placeholders module.
"""

# Path setup handled by conftest.py

from pptx import Presentation
from pptx_gen.master_placeholders import add_master_placeholders


def test_footer_placeholder_none_handling():
    """Test that None placeholders are handled gracefully."""
    prs = Presentation.create(template="mitsein")
    master = prs.slide_masters[0]

    # Test with page-number enabled
    placeholders = {
        'page-number': {
            'enabled': True,
            'position': 'br',
            'positionConfig': {'bottom': 20, 'right': 40}
        }
    }

    theme = {'accent': '#FF0000', 'text_muted': '#888888'}

    # Should not crash even if placeholder is None
    add_master_placeholders(master, placeholders, theme)

    print("✅ Footer placeholder None check test passed")


def test_date_placeholder_none_handling():
    """Test that None date placeholder is handled gracefully."""
    prs = Presentation.create(template="mitsein")
    master = prs.slide_masters[0]

    # Test with date enabled
    placeholders = {
        'date': {
            'enabled': True,
            'position': 'bl',
            'positionConfig': {'bottom': 20, 'left': 40}
        }
    }

    theme = {'accent': '#FF0000', 'text_muted': '#888888'}

    # Should not crash even if placeholder is None
    add_master_placeholders(master, placeholders, theme)

    print("✅ Date placeholder None check test passed")


if __name__ == '__main__':
    test_footer_placeholder_none_handling()
    test_date_placeholder_none_handling()
    print("\n✅ All master_placeholders tests passed!")
