"""
PPTX Generator - Generate PowerPoint from Layout Generator JSON config.

Uses python-pptx Slide Master API for proper inheritance.
"""

import sys
import json
import os
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Add python-pptx (enhanced fork) to path
# Try to use environment variable, fallback to default location
PPTX_LIB_PATH = os.environ.get('PYTHON_PPTX_PATH', '/Users/touichskyer/Code/python-pptx/src')
sys.path.insert(0, PPTX_LIB_PATH)

from pptx import Presentation

from .themes import THEME_COLORS, get_theme
from .slide_master import setup_slide_master
from .adapter import normalize_input, to_legacy_config_for_page
from .slides import (
    generate_cover_slide,
    generate_divider_slide,
    generate_smartart_slide,
    generate_grid_slide,
)


def generate_pptx(config: dict, output_path: str):
    """Generate PPTX from config using Slide Master API.

    Args:
        config: Configuration dict with theme, slideMaster, grid, etc.
        output_path: Path to save the generated PPTX file.
    """
    doc = normalize_input(config)
    master = doc.get('master', {})
    pages = doc.get('pages', [])
    theme_id = master.get('theme', 'soft_peach_cream')
    theme = get_theme(theme_id)

    # Build theme dict for Presentation.create() API
    accent_colors = theme.get('accentColors', [theme['accent']] * 6)
    pptx_theme = {
        'dk1': theme['primary'],
        'lt1': theme['bg'],
        'accent1': accent_colors[0] if len(accent_colors) > 0 else theme['accent'],
        'accent2': accent_colors[1] if len(accent_colors) > 1 else theme['accent'],
        'accent3': accent_colors[2] if len(accent_colors) > 2 else theme['accent'],
        'accent4': accent_colors[3] if len(accent_colors) > 3 else theme['accent'],
        'accent5': accent_colors[4] if len(accent_colors) > 4 else theme['accent'],
        'accent6': accent_colors[5] if len(accent_colors) > 5 else theme['accent'],
        'major_font': theme.get('major_font', 'Arial'),
        'minor_font': theme.get('minor_font', 'Arial')
    }

    # Create presentation with mitsein template (has 18 layouts) and theme applied
    prs = Presentation.create(template="mitsein", theme=pptx_theme)

    # Setup slide master once for the whole presentation.
    master_config = {
        'theme': theme_id,
        'slide': doc.get('slide', {}),
        'slideMaster': {
            'decorativeShapes': master.get('masterShapes', []),
            'placeholders': master.get('masterPlaceholders', {}),
            'contentAreas': master.get('masterContentAreas', {}),
        },
    }
    setup_slide_master(prs, master_config)

    if not isinstance(pages, list) or len(pages) == 0:
        pages = [{'id': 'page-1', 'type': 'content-grid', 'data': {'grid': {'layout': 'single', 'zones': []}}}]

    for page in pages:
        page_config = to_legacy_config_for_page(doc, page)
        page_type = page_config.get('pageType', 'content-grid')

        if page_type == 'cover':
            generate_cover_slide(prs, page_config, theme)
        elif page_type == 'divider':
            generate_divider_slide(prs, page_config, theme)
        elif page_type == 'content-smartart':
            generate_smartart_slide(prs, page_config, theme)
        else:
            generate_grid_slide(prs, page_config, theme)

    prs.save(output_path)
    logger.info(f"Saved PPTX: {output_path}")


def main():
    """CLI entry point."""
    if len(sys.argv) < 2:
        demo_config = {
            "theme": "soft_peach_cream",
            "pageType": "content-grid",
            "slideMaster": {
                "decorativeShapes": [
                    {
                        "id": "side-bar",
                        "configType": "thickness-positions",
                        "thickness": "thin",
                        "positions": ["left"],
                        "thicknessConfig": {"size": 8, "boundsOffset": 20}
                    },
                    {
                        "id": "corner-dots",
                        "configType": "thickness-positions",
                        "thickness": "small",
                        "positions": ["br"],
                        "thicknessConfig": {"size": 40, "dotSize": 2, "gap": 8},
                        "style": {"opacity": 0.3}
                    }
                ],
                "placeholders": {
                    "page-number": {
                        "enabled": True,
                        "position": "br",
                        "positionConfig": {"bottom": 20, "right": 60}
                    },
                    "logo": {
                        "enabled": True,
                        "position": "tl",
                        "positionConfig": {"top": 24, "left": 40},
                        "size": "medium",
                        "sizeConfig": {"width": 80, "height": 30}
                    }
                },
                "contentAreas": {
                    "titleStyle": "with-tag",
                    "sourceStyle": "citation",
                    "headerBounds": {"x": 40, "y": 20, "width": 1200, "height": 60},
                    "bodyBounds": {"x": 40, "y": 88, "width": 1200, "height": 552},
                    "footerBounds": {"x": 40, "y": 680, "width": 1200, "height": 40}
                }
            },
            "grid": {
                "layout": "two-col-equal",
                "zones": [
                    {"id": "A", "content": "chart", "flex": 1},
                    {"id": "B", "content": "text", "flex": 1}
                ]
            }
        }
        # Use current directory for demo output
        output_path = os.path.join(os.getcwd(), "output.pptx")
        generate_pptx(demo_config, output_path)
        logger.info(f"Demo PPTX generated: {output_path}")
        return

    config_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"

    with open(config_path, 'r', encoding='utf-8') as f:
        config = json.load(f)

    generate_pptx(config, output_path)


__all__ = [
    'generate_pptx',
    'main',
    'THEME_COLORS',
    'get_theme',
    'setup_slide_master',
    'generate_cover_slide',
    'generate_divider_slide',
    'generate_smartart_slide',
    'generate_grid_slide',
]
