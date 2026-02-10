"""
Unit tests for PPTX theme generation.

Tests that theme colors (accent1-6, dk1, lt1, etc.) are correctly
applied when generating PPTX files.
"""

import os
import tempfile

# Path setup handled by conftest.py

import pytest
from pptx import Presentation

from pptx_gen import generate_pptx
from pptx_gen.themes import THEME_COLORS, get_theme


class TestThemeColors:
    """Test that theme colors are correctly defined."""

    def test_all_themes_have_accent_colors(self):
        """All themes should have accentColors array with 6 colors."""
        for theme_id, theme in THEME_COLORS.items():
            assert 'accentColors' in theme, f"Theme '{theme_id}' missing accentColors"
            assert len(theme['accentColors']) == 6, f"Theme '{theme_id}' should have 6 accent colors"

    def test_all_themes_have_required_fields(self):
        """All themes should have required color fields."""
        required_fields = ['primary', 'accent', 'bg', 'text', 'text_muted']
        for theme_id, theme in THEME_COLORS.items():
            for field in required_fields:
                assert field in theme, f"Theme '{theme_id}' missing field '{field}'"

    def test_get_theme_returns_default_for_unknown(self):
        """get_theme should return default theme for unknown theme_id."""
        theme = get_theme('nonexistent_theme')
        assert theme is not None
        assert theme == THEME_COLORS['soft_peach_cream']

    def test_accent_colors_are_valid_hex(self):
        """All accent colors should be valid hex strings."""
        import re
        hex_pattern = re.compile(r'^#[0-9A-Fa-f]{6}$')

        for theme_id, theme in THEME_COLORS.items():
            for i, color in enumerate(theme['accentColors']):
                assert hex_pattern.match(color), \
                    f"Theme '{theme_id}' accent color {i+1} '{color}' is not valid hex"


class TestPresentationCreate:
    """Test Presentation.create() with theme parameter."""

    def test_create_with_theme_sets_colors(self):
        """Presentation.create(theme=...) should set all color scheme colors."""
        theme = {
            'dk1': '#1A2942',
            'lt1': '#FFFFFF',
            'accent1': '#FF0000',
            'accent2': '#00FF00',
            'accent3': '#0000FF',
            'accent4': '#FFFF00',
            'accent5': '#FF00FF',
            'accent6': '#00FFFF',
            'major_font': 'Arial',
            'minor_font': 'Calibri',
        }

        prs = Presentation.create(theme=theme)
        color_scheme = prs.slide_masters[0].theme.color_scheme

        assert str(color_scheme.dk1) == '1A2942'
        assert str(color_scheme.lt1) == 'FFFFFF'
        assert str(color_scheme.accent1) == 'FF0000'
        assert str(color_scheme.accent2) == '00FF00'
        assert str(color_scheme.accent3) == '0000FF'
        assert str(color_scheme.accent4) == 'FFFF00'
        assert str(color_scheme.accent5) == 'FF00FF'
        assert str(color_scheme.accent6) == '00FFFF'

    def test_create_derives_missing_colors(self):
        """Presentation.create() should derive dk2, lt2, hlink, folHlink."""
        theme = {
            'dk1': '#000000',
            'lt1': '#FFFFFF',
            'accent1': '#FF0000',
            'major_font': 'Arial',
            'minor_font': 'Calibri',
        }

        prs = Presentation.create(theme=theme)
        color_scheme = prs.slide_masters[0].theme.color_scheme

        # These should be auto-derived
        assert color_scheme.dk2 is not None
        assert color_scheme.lt2 is not None
        assert color_scheme.hlink is not None
        assert color_scheme.folHlink is not None


class TestGeneratePptx:
    """Test the generate_pptx function."""

    def test_generate_with_soft_peach_cream(self):
        """Generate PPTX with soft_peach_cream theme."""
        config = {
            'theme': 'soft_peach_cream',
            'pageType': 'content-grid',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'grid': {'layout': 'single', 'zones': [{'id': 'A', 'content': 'text', 'flex': 1}]},
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)

            # Verify the file was created and can be opened
            prs = Presentation.open(output_path)
            color_scheme = prs.slide_masters[0].theme.color_scheme

            # Check that accent colors from soft_peach_cream are applied
            theme = get_theme('soft_peach_cream')
            expected_accent1 = theme['accentColors'][0].lstrip('#').upper()
            assert str(color_scheme.accent1) == expected_accent1
        finally:
            os.unlink(output_path)

    def test_generate_with_executive_theme(self):
        """Generate PPTX with executive (dark) theme."""
        config = {
            'theme': 'executive',
            'pageType': 'content-grid',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'grid': {'layout': 'single', 'zones': [{'id': 'A', 'content': 'text', 'flex': 1}]},
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)

            prs = Presentation.open(output_path)
            color_scheme = prs.slide_masters[0].theme.color_scheme

            theme = get_theme('executive')
            expected_accent1 = theme['accentColors'][0].lstrip('#').upper()
            assert str(color_scheme.accent1) == expected_accent1
        finally:
            os.unlink(output_path)

    def test_generate_with_premium_report_theme(self):
        """Generate PPTX with premium report theme (deep_red_blue)."""
        config = {
            'theme': 'deep_red_blue',
            'pageType': 'content-grid',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'grid': {'layout': 'single', 'zones': [{'id': 'A', 'content': 'text', 'flex': 1}]},
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)

            prs = Presentation.open(output_path)
            color_scheme = prs.slide_masters[0].theme.color_scheme

            theme = get_theme('deep_red_blue')
            # deep_red_blue accent1 should be #731817
            assert str(color_scheme.accent1) == '731817'
        finally:
            os.unlink(output_path)


class TestColorSchemeAPI:
    """Test the ColorScheme API directly."""

    def test_color_scheme_accepts_hex_string(self):
        """ColorScheme should accept hex strings with or without #."""
        prs = Presentation.create()
        color_scheme = prs.slide_masters[0].theme.color_scheme

        # With #
        color_scheme.accent1 = '#FF5500'
        assert str(color_scheme.accent1) == 'FF5500'

        # Without #
        color_scheme.accent2 = '00FF55'
        assert str(color_scheme.accent2) == '00FF55'

    def test_color_scheme_iteration(self):
        """ColorScheme should be iterable."""
        prs = Presentation.create()
        color_scheme = prs.slide_masters[0].theme.color_scheme

        colors = dict(color_scheme)
        assert 'dk1' in colors
        assert 'lt1' in colors
        assert 'accent1' in colors
        assert 'accent6' in colors


if __name__ == '__main__':
    pytest.main([__file__, '-v'])
