"""
Unit tests for SmartArt generation.

Tests that SmartArt color schemes work correctly with theme colors.
"""

import os
import tempfile
import zipfile

# Path setup handled by conftest.py

import pytest
from pptx import Presentation
from pptx.enum.smartart import SMARTART_TYPE, SMARTART_COLORS

from pptx_gen import generate_pptx
from pptx_gen.themes import THEME_COLORS, get_theme
from pptx_gen.slides.smartart import (
    _create_smartart_data,
    _extract_smartart_items,
    _resolve_smartart_font_size_pt,
)


class TestSmartArtColorSchemes:
    """Test SmartArt color scheme mappings."""

    def test_smartart_colors_enum_exists(self):
        """SMARTART_COLORS enum should have colorful schemes."""
        assert hasattr(SMARTART_COLORS, 'COLORFUL_ACCENT_COLORS')
        assert hasattr(SMARTART_COLORS, 'COLORFUL_ACCENT_2_TO_3')
        assert hasattr(SMARTART_COLORS, 'ACCENT_1')
        assert hasattr(SMARTART_COLORS, 'ACCENT_2')

    def test_smartart_type_enum_exists(self):
        """SMARTART_TYPE enum should have pyramid, matrix, etc."""
        assert hasattr(SMARTART_TYPE, 'BASIC_PYRAMID')
        assert hasattr(SMARTART_TYPE, 'BASIC_MATRIX')
        assert hasattr(SMARTART_TYPE, 'BASIC_CYCLE')
        assert hasattr(SMARTART_TYPE, 'BASIC_PROCESS')


class TestSmartArtGeneration:
    """Test SmartArt generation with different themes."""

    def test_generate_smartart_slide(self):
        """Generate PPTX with SmartArt slide."""
        config = {
            'theme': 'soft_peach_cream',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'smartart': {
                'type': 'pyramid',
                'category': 'pyramid',
                'placement': 'full',
                'colorScheme': 'colorful1',
            },
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)

            # Verify file was created
            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 0

            # Can open the file
            prs = Presentation.open(output_path)
            assert len(prs.slides) == 1
        finally:
            os.unlink(output_path)

    def test_smartart_with_premium_theme(self):
        """Generate SmartArt with premium report theme."""
        config = {
            'theme': 'deep_blue_gold',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'smartart': {
                'type': 'chevron',
                'category': 'process',
                'placement': 'full',
                'colorScheme': 'colorful1',
            },
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)

            prs = Presentation.open(output_path)
            color_scheme = prs.slide_masters[0].theme.color_scheme

            # Verify theme colors are set
            theme = get_theme('deep_blue_gold')
            expected_accent1 = theme['accentColors'][0].lstrip('#').upper()
            assert str(color_scheme.accent1) == expected_accent1
        finally:
            os.unlink(output_path)

    def test_smartart_uses_content_shell_title_and_footer(self):
        """SmartArt slide should honor content-shell header/footer fields."""
        config = {
            'theme': 'forest_green',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {
                'decorativeShapes': [],
                'placeholders': {},
                'contentAreas': {
                    'titleStyle': 'with-tag',
                    'sourceStyle': 'citation',
                    'headerBounds': {'x': 40, 'y': 20, 'width': 1200, 'height': 60},
                    'bodyBounds': {'x': 40, 'y': 88, 'width': 1200, 'height': 592},
                    'footerBounds': {'x': 40, 'y': 680, 'width': 1200, 'height': 24},
                },
            },
            'smartart': {
                'type': 'pyramid',
                'category': 'pyramid',
                'placement': 'left-desc',
                'colorScheme': 'colorful1',
            },
            'contentTitle': '自定义标题',
            'contentTag': '专题',
            'contentSource': '内部系统',
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)
            prs = Presentation.open(output_path)
            slide = prs.slides[0]
            text_values = []
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                text = shape.text_frame.text.strip()
                if text:
                    text_values.append(text)

            joined = '\n'.join(text_values)
            assert '自定义标题' in joined
            assert '专题' in joined
            assert '数据来源：内部系统' in joined
        finally:
            os.unlink(output_path)

    def test_smartart_font_size_pt_is_written_to_data_xml(self):
        """Explicit smartart.fontSizePt should be persisted into SmartArt data.xml runs."""
        config = {
            'theme': 'forest_green',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'smartart': {
                'type': 'pyramid',
                'category': 'pyramid',
                'placement': 'full',
                'colorScheme': 'colorful1',
                'fontSizePt': 11,
                'items': ['A', 'B', 'C'],
            },
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)
            with zipfile.ZipFile(output_path, 'r') as zf:
                data_parts = [name for name in zf.namelist() if name.startswith('ppt/diagrams/data') and name.endswith('.xml')]
                assert data_parts
                xml = zf.read(data_parts[0]).decode('utf-8')
                assert 'a:rPr sz="1100"' in xml
        finally:
            os.unlink(output_path)


class TestSmartArtFontSizeResolve:
    """Test SmartArt font size resolution and clamping."""

    def test_resolve_font_size_uses_explicit_config(self):
        assert _resolve_smartart_font_size_pt({'fontSizePt': 10.5}, 'pyramid') == 10.5

    def test_resolve_font_size_uses_type_default(self):
        assert _resolve_smartart_font_size_pt({}, 'square-accent-list') == 10.0
        assert _resolve_smartart_font_size_pt({}, 'matrix') == 12.0

    def test_resolve_font_size_clamps_bounds(self):
        assert _resolve_smartart_font_size_pt({'fontSizePt': 1}, 'pyramid') == 8.0
        assert _resolve_smartart_font_size_pt({'fontSizePt': 99}, 'pyramid') == 36.0


class TestSmartArtAPI:
    """Test python-pptx SmartArt API directly."""

    def test_add_smartart_to_slide(self):
        """Can add SmartArt to a slide."""
        from pptx.util import Inches
        from pptx.smartart import SmartArtData

        prs = Presentation.create()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        data = SmartArtData()
        data.add_node('Item 1')
        data.add_node('Item 2')
        data.add_node('Item 3')

        smartart = slide.shapes.add_smartart(
            smartart_type=SMARTART_TYPE.BASIC_PYRAMID,
            x=Inches(1),
            y=Inches(1),
            cx=Inches(8),
            cy=Inches(5),
            smartart_data=data,
            color_scheme=SMARTART_COLORS.COLORFUL_ACCENT_COLORS,
        )

        assert smartart is not None

    def test_smartart_color_scheme_options(self):
        """SmartArt should support different color schemes."""
        from pptx.util import Inches
        from pptx.smartart import SmartArtData

        prs = Presentation.create()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Test different color schemes
        schemes = [
            SMARTART_COLORS.COLORFUL_ACCENT_COLORS,
            SMARTART_COLORS.COLORFUL_ACCENT_2_TO_3,
            SMARTART_COLORS.ACCENT_1,
            SMARTART_COLORS.ACCENT_2,
        ]

        for i, scheme in enumerate(schemes):
            data = SmartArtData()
            data.add_node('A')
            data.add_node('B')

            smartart = slide.shapes.add_smartart(
                smartart_type=SMARTART_TYPE.BASIC_PROCESS,
                x=Inches(i * 2),
                y=Inches(1),
                cx=Inches(1.5),
                cy=Inches(1),
                smartart_data=data,
                color_scheme=scheme,
            )
            assert smartart is not None


class TestThemeAccentColorsInSmartArt:
    """Test that theme accent colors are used by SmartArt."""

    def test_custom_accent_colors_applied(self):
        """Custom accent colors should be reflected in the theme."""
        theme = {
            'dk1': '#000000',
            'lt1': '#FFFFFF',
            'accent1': '#FF0000',  # Red
            'accent2': '#00FF00',  # Green
            'accent3': '#0000FF',  # Blue
            'accent4': '#FFFF00',  # Yellow
            'accent5': '#FF00FF',  # Magenta
            'accent6': '#00FFFF',  # Cyan
            'major_font': 'Arial',
            'minor_font': 'Arial',
        }

        prs = Presentation.create(theme=theme)
        color_scheme = prs.slide_masters[0].theme.color_scheme

        # Verify all accent colors
        assert str(color_scheme.accent1) == 'FF0000'
        assert str(color_scheme.accent2) == '00FF00'
        assert str(color_scheme.accent3) == '0000FF'
        assert str(color_scheme.accent4) == 'FFFF00'
        assert str(color_scheme.accent5) == 'FF00FF'
        assert str(color_scheme.accent6) == '00FFFF'


class TestSmartArtLayouts:
    """Test specific SmartArt layout types."""

    def test_cycle_layout(self):
        """Generate PPTX with basic cycle (cycle4) SmartArt."""
        config = {
            'theme': 'soft_peach_cream',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'smartart': {
                'type': 'cycle',
                'category': 'cycle',
                'placement': 'full',
                'colorScheme': 'colorful1',
                'items': ['战略愿景', '日常任务', '经营目标', '执行计划'],
            },
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)
            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 0

            prs = Presentation.open(output_path)
            assert len(prs.slides) == 1
        finally:
            os.unlink(output_path)


    def test_hexagon_alternating_layout(self):
        """Generate PPTX with alternating hexagons SmartArt."""
        config = {
            'theme': 'deep_blue_gold',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'smartart': {
                'type': 'hexagon-alternating',
                'category': 'list',
                'placement': 'full',
                'colorScheme': 'colorful1',
                'items': ['Item 1', 'Item 2', 'Item 3', 'Item 4', 'Item 5', 'Item 6'],
            },
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)
            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 0

            prs = Presentation.open(output_path)
            assert len(prs.slides) == 1
        finally:
            os.unlink(output_path)

    def test_pyramid_list_layout(self):
        """Generate PPTX with pyramid list (pyramid2) SmartArt."""
        config = {
            'theme': 'soft_peach_cream',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'smartart': {
                'type': 'pyramid-list',
                'category': 'pyramid',
                'placement': 'full',
                'colorScheme': 'colorful1',
                'items': ['Level 1', 'Level 2', 'Level 3', 'Level 4', 'Level 5'],
            },
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)
            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 0

            prs = Presentation.open(output_path)
            assert len(prs.slides) == 1
        finally:
            os.unlink(output_path)

    @pytest.mark.parametrize(
        ("type_id", "layout_suffix"),
        [
            ("cycle1", "cycle1"),
            ("cycle2", "cycle2"),
            ("cycle3", "cycle3"),
            ("cycle5", "cycle5"),
            ("cycle6", "cycle6"),
            ("cycle7", "cycle7"),
            ("cycle8", "cycle8"),
        ],
    )
    def test_cycle_ooxml_layout_mapping(self, type_id, layout_suffix):
        """Cycle 9-15 mappings should export expected OOXML layout URIs."""
        config = {
            'theme': 'forest_green',
            'pageType': 'content-smartart',
            'slide': {'width': 1280, 'height': 720, 'widthInches': 13.333, 'heightInches': 7.5},
            'slideMaster': {'decorativeShapes': [], 'placeholders': {}, 'contentAreas': {}},
            'smartart': {
                'type': type_id,
                'category': 'cycle',
                'placement': 'full',
                'colorScheme': 'colorful1',
                'items': ['A', 'B', 'C', 'D', 'E'],
            },
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            generate_pptx(config, output_path)
            with zipfile.ZipFile(output_path) as zf:
                rel_xml = zf.read('ppt/slides/_rels/slide1.xml.rels').decode('utf-8')
                import re

                match = re.search(r'Target="../diagrams/(layout\d+\.xml)"', rel_xml)
                assert match, "diagram layout relationship not found"
                layout_xml = zf.read(f"ppt/diagrams/{match.group(1)}").decode('utf-8')
                unique = re.search(r'uniqueId="([^"]+)"', layout_xml)
                assert unique, "layout uniqueId not found"
                assert unique.group(1).endswith(f'/{layout_suffix}')
        finally:
            os.unlink(output_path)


class TestSmartArtContract:
    """Test SmartArt contract field handling."""

    def test_extract_items_prefers_explicit_items(self):
        """smartart.items should take precedence over smartart.ooxml.items."""
        config = {
            'items': ['A', {'text': 'B'}],
            'ooxml': {'items': [{'text': 'legacy-1'}, {'text': 'legacy-2'}]},
        }
        assert _extract_smartart_items(config) == ['A', {'text': 'B'}]

    def test_extract_items_falls_back_to_ooxml(self):
        """Fallback to ooxml.items when explicit items is absent."""
        config = {
            'ooxml': {'items': [{'text': 'X'}, 'Y']},
        }
        assert _extract_smartart_items(config) == [{'text': 'X'}, 'Y']

    def test_hierarchy_children_are_preserved(self):
        """Hierarchy nodes should keep nested children from structured items."""
        items = [
            {
                'text': 'CEO',
                'children': [
                    {'text': 'CTO'},
                    {'text': 'CFO', 'children': [{'text': 'Finance'}]},
                ],
            }
        ]

        data = _create_smartart_data('hierarchy', items)
        assert len(data.root_nodes) == 1
        assert data.root_nodes[0].text == 'CEO'
        assert [n.text for n in data.root_nodes[0].children] == ['CTO', 'CFO']
        assert data.root_nodes[0].children[1].children[0].text == 'Finance'

    def test_pyramid_children_are_preserved(self):
        """Pyramid nodes should keep child bullet text when provided."""
        items = [
            {'text': '愿景', 'children': [{'text': '长期目标'}]},
            {'text': '战略', 'children': [{'text': '中期规划'}]},
        ]

        data = _create_smartart_data('pyramid', items)
        assert len(data.root_nodes) == 2
        assert data.root_nodes[0].text == '愿景'
        assert [n.text for n in data.root_nodes[0].children] == ['长期目标']
        assert [n.text for n in data.root_nodes[1].children] == ['中期规划']

    def test_cycle_children_are_preserved(self):
        """Cycle nodes should keep child description text when provided."""
        items = [
            {'text': '计划', 'children': [{'text': 'Plan'}]},
            {'text': '执行', 'children': [{'text': 'Do'}]},
        ]

        data = _create_smartart_data('cycle', items)
        assert len(data.root_nodes) == 2
        assert [n.text for n in data.root_nodes[0].children] == ['Plan']
        assert [n.text for n in data.root_nodes[1].children] == ['Do']

    def test_list_variant_children_are_preserved(self):
        """List-like variants should also preserve nested child bullets."""
        items = [
            {'text': '要点一', 'children': [{'text': '子项A'}, {'text': '子项B'}]},
            {'text': '要点二', 'children': [{'text': '子项C'}]},
        ]

        data = _create_smartart_data('square-accent-list', items)
        assert len(data.root_nodes) == 2
        assert [n.text for n in data.root_nodes[0].children] == ['子项A', '子项B']
        assert [n.text for n in data.root_nodes[1].children] == ['子项C']


if __name__ == '__main__':
    pytest.main([__file__, '-v'])
