"""Tests for multi-page payload support and v1 backward compatibility."""

from __future__ import annotations

import os
import tempfile

from pptx import Presentation

from pptx_gen import generate_pptx
from pptx_gen.adapter import normalize_input, to_legacy_config_for_page


def test_slide_count_equals_pages_count_for_v2():
    config = {
        "schemaVersion": 2,
        "master": {
            "theme": "azure",
            "masterShapes": [],
            "masterPlaceholders": {},
            "masterContentAreas": {},
        },
        "pages": [
            {"id": "p1", "type": "cover", "data": {"coverLayout": "cross_rectangles"}},
            {"id": "p2", "type": "divider", "data": {"divider": {"layout": "cards-highlight"}}},
            {"id": "p3", "type": "content-grid", "data": {"grid": {"layout": "single", "zones": []}}},
        ],
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fp:
        output_path = fp.name

    try:
        generate_pptx(config, output_path)
        prs = Presentation.open(output_path)
        assert len(prs.slides) == 3
    finally:
        os.unlink(output_path)


def test_v1_backward_compat_still_generates_single_page():
    config = {
        "theme": "azure",
        "pageType": "cover",
        "slideMaster": {"decorativeShapes": [], "placeholders": {}, "contentAreas": {}},
        "coverLayout": "triangle_stack",
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fp:
        output_path = fp.name

    try:
        generate_pptx(config, output_path)
        prs = Presentation.open(output_path)
        assert len(prs.slides) == 1
    finally:
        os.unlink(output_path)


def test_normalize_v1_to_v2_uses_slide_master_fields():
    v1 = {
        "theme": "forest_green",
        "pageType": "content-smartart",
        "slideMaster": {
            "decorativeShapes": [{"id": "side-bar"}],
            "placeholders": {"logo": {"enabled": True}},
            "contentAreas": {"titleStyle": "with-tag"},
        },
        "smartart": {"type": "pyramid", "items": [{"text": "A"}]},
    }

    v2 = normalize_input(v1)
    assert v2["schemaVersion"] == 2
    assert v2["master"]["theme"] == "forest_green"
    assert v2["master"]["masterShapes"] == [{"id": "side-bar"}]
    assert v2["master"]["masterPlaceholders"] == {"logo": {"enabled": True}}
    assert v2["master"]["masterContentAreas"] == {"titleStyle": "with-tag"}
    assert len(v2["pages"]) == 1
    assert v2["pages"][0]["type"] == "content-smartart"
    assert v2["pages"][0]["shell"] == "content"
    assert v2["pages"][0]["renderer"] == "smartart"
    assert v2["pages"][0]["pageShell"] == "content"
    assert v2["pages"][0]["bodyRenderer"] == "smartart"


def test_divider_flat_fields_override_nested_divider_payload():
    doc = normalize_input(
        {
            "schemaVersion": 2,
            "master": {
                "theme": "forest_green",
                "masterShapes": [],
                "masterPlaceholders": {},
                "masterContentAreas": {},
            },
            "pages": [
                {
                    "id": "p1",
                    "type": "divider",
                    "data": {
                        "divider": {
                            "layout": "cards",
                            "sectionIndex": 2,
                            "bgStyle": "light",
                        },
                        "dividerLayout": "cards-highlight",
                        "dividerIndex": 1,
                        "dividerBgStyle": "solid",
                    },
                }
            ],
        }
    )

    legacy = to_legacy_config_for_page(doc, doc["pages"][0])
    divider = legacy["divider"]
    assert divider["layout"] == "cards-highlight"
    assert divider["sectionIndex"] == 1
    assert divider["bgStyle"] == "solid"


def test_normalize_v2_supports_page_shell_body_renderer_without_type():
    v2 = normalize_input(
        {
            "schemaVersion": 2,
            "master": {
                "theme": "forest_green",
                "masterShapes": [],
                "masterPlaceholders": {},
                "masterContentAreas": {},
            },
            "pages": [
                {
                    "id": "p1",
                    "pageShell": "content",
                    "bodyRenderer": "smartart",
                    "bodyLayout": "left-desc",
                    "data": {
                        "smartartType": "pyramid",
                        "smartartItemsByType": {"pyramid": [{"text": "A"}]},
                    },
                }
            ],
        }
    )

    page = v2["pages"][0]
    assert page["type"] == "content-smartart"
    assert page["shell"] == "content"
    assert page["renderer"] == "smartart"
    assert page["layout"] == "left-desc"


def test_to_legacy_config_resolves_type_from_page_model():
    doc = normalize_input(
        {
            "schemaVersion": 2,
            "master": {
                "theme": "forest_green",
                "masterShapes": [],
                "masterPlaceholders": {},
                "masterContentAreas": {},
            },
            "pages": [
                {
                    "id": "p1",
                    "shell": "content",
                    "renderer": "grid",
                    "layout": "single",
                    "data": {"grid": {"layout": "single", "zones": []}},
                }
            ],
        }
    )

    legacy = to_legacy_config_for_page(doc, doc["pages"][0])
    assert legacy["pageType"] == "content-grid"


def test_to_legacy_config_carries_content_shell_fields():
    doc = normalize_input(
        {
            "schemaVersion": 2,
            "master": {
                "theme": "forest_green",
                "masterShapes": [],
                "masterPlaceholders": {},
                "masterContentAreas": {},
            },
            "pages": [
                {
                    "id": "p1",
                    "type": "content-smartart",
                    "data": {
                        "contentTitle": "自定义标题",
                        "contentTag": "专题",
                        "contentSource": "内部系统",
                        "smartartType": "pyramid",
                    },
                }
            ],
        }
    )

    legacy = to_legacy_config_for_page(doc, doc["pages"][0])
    assert legacy["contentTitle"] == "自定义标题"
    assert legacy["contentTag"] == "专题"
    assert legacy["contentSource"] == "内部系统"
