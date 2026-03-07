"""
Divider/section slide generator (aligned with frontend divider model).
"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..themes import hex_to_rgb
from ..utils import get_blank_layout, remove_slide_placeholders


SECTION_DATA = [
    {"num": "1", "zh": "年度工作概述", "en": "ANNUAL WORK OVERVIEW"},
    {"num": "2", "zh": "工作完成情况", "en": "WORK COMPLETION"},
    {"num": "3", "zh": "项目成果展示", "en": "PROJECT RESULTS"},
    {"num": "4", "zh": "工作不足与改进", "en": "IMPROVEMENTS"},
    {"num": "5", "zh": "未来发展规划", "en": "FUTURE PLANS"},
    {"num": "6", "zh": "总结与展望", "en": "SUMMARY & OUTLOOK"},
]

ROMAN_NUMERALS = ["I", "II", "III", "IV", "V", "VI"]
CHINESE_NUMERALS = ["一", "二", "三", "四", "五", "六"]
CIRCLED_NUMERALS = ["①", "②", "③", "④", "⑤", "⑥"]


def generate_divider_slide(prs, config, theme):
    """Generate divider/section slide from frontend divider config."""
    slide = prs.slides.add_slide(get_blank_layout(prs))
    remove_slide_placeholders(slide)

    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400
    geom = _Geom(slide_w, slide_h)

    divider = config.get("divider", {}) if isinstance(config.get("divider"), dict) else {}
    layout = str(divider.get("layout", "cards-highlight"))
    section_count = int(divider.get("sectionCount", 4) or 4)
    section_count = max(1, min(6, section_count))
    number_style = str(divider.get("numberStyle", "arabic"))
    text_level = str(divider.get("textLevel", "full"))
    bg_style = str(divider.get("bgStyle", "solid"))
    section_index = int(divider.get("sectionIndex", 0) or 0)

    sections = SECTION_DATA[:section_count]
    _apply_background(slide, theme, bg_style)
    text_color = _text_color_for_bg(theme, bg_style)

    if layout in {"cards-highlight", "cards", "strips"}:
        _render_toc_layout(
            slide=slide,
            layout=layout,
            sections=sections,
            section_index=section_index,
            number_style=number_style,
            compact=(text_level == "compact"),
            theme=theme,
            text_color=text_color,
            geom=geom,
        )
        return slide

    selected = _select_section(sections, section_index)
    _render_section_layout(
        slide=slide,
        layout=layout,
        section=selected,
        number_style=number_style,
        compact=(text_level == "compact"),
        theme=theme,
        text_color=text_color,
        geom=geom,
    )
    return slide


def _apply_background(slide, theme: dict, bg_style: str) -> None:
    accent = _rgb(theme.get("accent", "#2D5A3D"))
    bg = _rgb(theme.get("bg", "#F2F5F6"))

    slide.background.fill.solid()
    if bg_style == "light":
        slide.background.fill.fore_color.rgb = bg
        return
    if bg_style == "split":
        slide.background.fill.fore_color.rgb = _shade(accent, 0.86)
        return
    if bg_style == "gradient":
        slide.background.fill.fore_color.rgb = _shade(accent, 0.78)
        return
    slide.background.fill.fore_color.rgb = accent


def _text_color_for_bg(theme: dict, bg_style: str) -> RGBColor:
    if bg_style == "light":
        return _rgb(theme.get("primary", "#101010"))
    accent = _rgb(theme.get("accent", "#2D5A3D"))
    if _is_dark(accent):
        return RGBColor(255, 255, 255)
    return _rgb(theme.get("primary", "#101010"))


def _render_toc_layout(
    *,
    slide,
    layout: str,
    sections: list[dict],
    section_index: int,
    number_style: str,
    compact: bool,
    theme: dict,
    text_color: RGBColor,
    geom,
) -> None:
    title = slide.shapes.add_textbox(geom.x(0.7), geom.y(0.45), geom.w(2.0), geom.h(0.6))
    p = title.text_frame.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = text_color

    if layout == "strips":
        _render_strips(slide, sections, section_index, number_style, compact, theme, text_color, geom)
        return

    if layout == "cards":
        _render_cards(
            slide, sections, section_index, number_style, compact, theme, text_color, geom, highlight=False
        )
        return

    _render_cards(
        slide, sections, section_index, number_style, compact, theme, text_color, geom, highlight=True
    )


def _render_cards(
    slide,
    sections: list[dict],
    section_index: int,
    number_style: str,
    compact: bool,
    theme: dict,
    text_color: RGBColor,
    geom,
    *,
    highlight: bool,
) -> None:
    left = 0.65
    right = 0.65
    top = 1.4
    height = 4.9
    gap = 0.14
    count = max(1, len(sections))
    card_w = (geom.base_w - left - right - gap * (count - 1)) / count
    active_idx = section_index - 1 if section_index > 0 else -1

    accent = _rgb(theme.get("accent", "#2D5A3D"))
    card_bg = _rgb(theme.get("card_bg", "#FFFFFF"))
    card_border = _rgb(theme.get("card_border", "#DADADA"))
    muted = _rgb(theme.get("text_muted", "#6F6F6F"))

    for i, section in enumerate(sections):
        x = left + i * (card_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            geom.x(x),
            geom.y(top),
            geom.w(card_w),
            geom.h(height),
        )
        rect.adjustments[0] = 0.03

        is_active = i == active_idx
        if highlight:
            rect.fill.solid()
            rect.fill.fore_color.rgb = accent if is_active else _shade(card_bg, 0.94)
            rect.line.fill.background()
        else:
            rect.fill.solid()
            rect.fill.fore_color.rgb = card_bg
            rect.line.color.rgb = accent if is_active else card_border
            rect.line.width = Pt(1.2 if is_active else 0.8)

        title_color = RGBColor(255, 255, 255) if (highlight and is_active) else text_color
        sub_color = RGBColor(240, 240, 240) if (highlight and is_active) else muted

        zh = slide.shapes.add_textbox(geom.x(x + 0.18), geom.y(top + 0.34), geom.w(card_w - 0.36), geom.h(0.9))
        p = zh.text_frame.paragraphs[0]
        p.text = section["zh"]
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = title_color

        if not compact:
            en = slide.shapes.add_textbox(geom.x(x + 0.18), geom.y(top + 1.05), geom.w(card_w - 0.36), geom.h(0.7))
            p = en.text_frame.paragraphs[0]
            p.text = section["en"]
            p.font.size = Pt(10)
            p.font.color.rgb = sub_color

        num = slide.shapes.add_textbox(
            geom.x(x + 0.12), geom.y(top + height - 1.25), geom.w(card_w - 0.2), geom.h(1.1)
        )
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        p = num.text_frame.paragraphs[0]
        p.text = _format_number(section["num"], number_style)
        p.font.size = Pt(74)
        if highlight and is_active:
            p.font.color.rgb = RGBColor(255, 255, 255)
        else:
            p.font.color.rgb = _alpha_like(title_color, 0.20)


def _render_strips(
    slide,
    sections: list[dict],
    section_index: int,
    number_style: str,
    compact: bool,
    theme: dict,
    text_color: RGBColor,
    geom,
) -> None:
    left_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, geom.x(0), geom.y(0), geom.w(2.1), geom.h(7.5))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = _shade(_rgb(theme.get("accent", "#2D5A3D")), 0.82)
    left_panel.line.fill.background()

    title = slide.shapes.add_textbox(geom.x(0.35), geom.y(2.0), geom.w(1.4), geom.h(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    right_left = 2.12
    strip_w = (geom.base_w - right_left) / max(1, len(sections))
    active_idx = section_index - 1 if section_index > 0 else -1
    accent = _rgb(theme.get("accent", "#2D5A3D"))

    for i, section in enumerate(sections):
        x = right_left + i * strip_w
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, geom.x(x), geom.y(0), geom.w(strip_w), geom.h(7.5))
        strip.fill.solid()
        if i == active_idx:
            strip.fill.fore_color.rgb = accent
        else:
            strip.fill.fore_color.rgb = _shade(accent, 0.92 - i * 0.05)
        strip.line.fill.background()

        text = slide.shapes.add_textbox(geom.x(x + 0.14), geom.y(1.15), geom.w(strip_w - 0.22), geom.h(3.2))
        p = text.text_frame.paragraphs[0]
        p.text = f"0{_format_number(section['num'], number_style)}"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p = text.text_frame.add_paragraph()
        p.text = section["zh"]
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(255, 255, 255)
        if not compact:
            p = text.text_frame.add_paragraph()
            p.text = section["en"]
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(240, 240, 240)


def _render_section_layout(
    *,
    slide,
    layout: str,
    section: dict,
    number_style: str,
    compact: bool,
    theme: dict,
    text_color: RGBColor,
    geom,
) -> None:
    part = _format_part(section["num"], number_style)
    number = _format_number(section["num"], number_style)
    accent = _rgb(theme.get("accent", "#2D5A3D"))

    if layout == "fullbleed":
        p = slide.shapes.add_textbox(geom.x(0.7), geom.y(0.9), geom.w(8.8), geom.h(1.0)).text_frame.paragraphs[0]
        p.text = part
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = text_color
    elif layout == "arrow":
        badge = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, geom.x(0.7), geom.y(1.0), geom.w(2.2), geom.h(1.0))
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        p = slide.shapes.add_textbox(geom.x(0.85), geom.y(1.28), geom.w(1.8), geom.h(0.5)).text_frame.paragraphs[0]
        p.text = part
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
    else:
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, geom.x(0.7), geom.y(0.78), geom.w(8.6), geom.h(0.02))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = _alpha_like(text_color, 0.25)
        top_line.line.fill.background()
        p = slide.shapes.add_textbox(geom.x(0.7), geom.y(1.05), geom.w(3.0), geom.h(0.6)).text_frame.paragraphs[0]
        p.text = part
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = text_color

    title = slide.shapes.add_textbox(geom.x(0.7), geom.y(2.05), geom.w(7.0), geom.h(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = section["zh"]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = text_color

    if not compact:
        subtitle = slide.shapes.add_textbox(geom.x(0.72), geom.y(3.25), geom.w(6.0), geom.h(0.6))
        p = subtitle.text_frame.paragraphs[0]
        p.text = section["en"]
        p.font.size = Pt(16)
        p.font.color.rgb = _alpha_like(text_color, 0.70)

    num = slide.shapes.add_textbox(geom.x(6.5), geom.y(1.0), geom.w(3.0), geom.h(5.0))
    if layout in {"left-align-mirror"}:
        num = slide.shapes.add_textbox(geom.x(0.1), geom.y(1.0), geom.w(3.0), geom.h(5.0))
    num.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    p = num.text_frame.paragraphs[0]
    p.text = number
    p.font.size = Pt(220)
    p.font.color.rgb = _alpha_like(text_color, 0.12)


def _format_number(raw_num: str, number_style: str) -> str:
    idx = max(0, int(raw_num) - 1)
    if number_style == "roman":
        return ROMAN_NUMERALS[idx] if idx < len(ROMAN_NUMERALS) else raw_num
    if number_style == "chinese":
        return CHINESE_NUMERALS[idx] if idx < len(CHINESE_NUMERALS) else raw_num
    if number_style == "circled":
        return CIRCLED_NUMERALS[idx] if idx < len(CIRCLED_NUMERALS) else raw_num
    return raw_num


def _format_part(raw_num: str, number_style: str) -> str:
    number = _format_number(raw_num, number_style)
    if number_style == "roman":
        return f"Part {number}"
    if number_style == "chinese":
        return f"第{number}部分"
    if number_style == "circled":
        return f"Part {number}"
    return f"第{number}部分"


def _select_section(sections: list[dict], section_index: int) -> dict:
    if not sections:
        return SECTION_DATA[0]
    if section_index > 0:
        idx = min(max(0, section_index - 1), len(sections) - 1)
        return sections[idx]
    return sections[0]


def _rgb(hex_color: str) -> RGBColor:
    return hex_to_rgb(hex_color)


def _shade(color: RGBColor, factor: float) -> RGBColor:
    factor = max(0.0, min(1.0, factor))
    return RGBColor(int(color[0] * factor), int(color[1] * factor), int(color[2] * factor))


def _alpha_like(color: RGBColor, alpha: float) -> RGBColor:
    alpha = max(0.0, min(1.0, alpha))
    return RGBColor(
        int(255 - (255 - color[0]) * alpha),
        int(255 - (255 - color[1]) * alpha),
        int(255 - (255 - color[2]) * alpha),
    )


def _is_dark(color: RGBColor) -> bool:
    y = 0.2126 * color[0] + 0.7152 * color[1] + 0.0722 * color[2]
    return y < 145


class _Geom:
    """Scale a 10x7.5 base coordinate system to current slide size."""

    def __init__(self, slide_w: float, slide_h: float):
        self.base_w = 10.0
        self.base_h = 7.5
        self.sx = slide_w / self.base_w if self.base_w else 1.0
        self.sy = slide_h / self.base_h if self.base_h else 1.0

    def x(self, value: float):
        return Inches(value * self.sx)

    def y(self, value: float):
        return Inches(value * self.sy)

    def w(self, value: float):
        return Inches(value * self.sx)

    def h(self, value: float):
        return Inches(value * self.sy)
